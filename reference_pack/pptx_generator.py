from __future__ import annotations

import re
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .content_builder import prepare_evidence_excerpt
from .evidence_renderer import EvidenceRenderer, EvidenceVisual
from .schemas import PreparedReference, ReferencePackRequest, SlideProvenance


TEXT = {
    "fr": {
        "section": "Références pertinentes\npour la mission",
        "summary": "Extrait de nos références {current}/{total}",
        "detail": "Référence détaillée {current}/{total}",
        "evidence": "Justificatifs de nos références {current}/{total}",
        "client": "Client",
        "country": "Pays",
        "period": "Période",
        "sector": "Secteur",
        "offering": "Offre / capacité",
        "description": "Description du projet",
        "services": "Services réalisés",
        "technologies": "Technologies et capacités",
        "why": "Pourquoi cette référence a été retenue",
        "source": "Source",
        "page": "page",
        "language": "langue",
        "prepared": "Préparé le",
        "fallback": "Text evidence fallback",
    },
    "en": {
        "section": "Relevant references\nfor the opportunity",
        "summary": "Selected references {current}/{total}",
        "detail": "Reference detail {current}/{total}",
        "evidence": "Reference evidence {current}/{total}",
        "client": "Client",
        "country": "Country",
        "period": "Period",
        "sector": "Sector",
        "offering": "Offering / capability",
        "description": "Project description",
        "services": "Services delivered",
        "technologies": "Technologies and capabilities",
        "why": "Why this reference was selected",
        "source": "Source",
        "page": "page",
        "language": "language",
        "prepared": "Prepared on",
        "fallback": "Text evidence fallback",
    },
    "ar": {
        "section": "المراجع ذات الصلة\nبالفرصة",
        "summary": "مختارات من مراجعنا {current}/{total}",
        "detail": "تفاصيل المرجع {current}/{total}",
        "evidence": "أدلة المراجع {current}/{total}",
        "client": "العميل",
        "country": "البلد",
        "period": "الفترة",
        "sector": "القطاع",
        "offering": "العرض / القدرة",
        "description": "وصف المشروع",
        "services": "الخدمات المنفذة",
        "technologies": "التقنيات والقدرات",
        "why": "سبب اختيار هذا المرجع",
        "source": "المصدر",
        "page": "الصفحة",
        "language": "اللغة",
        "prepared": "أُعد في",
        "fallback": "Text evidence fallback",
    },
}

PROHIBITED_TEXT = re.compile(
    r"\b(?:bm25|dense cosine|rrf|hybrid score|confidence percentage|chunk id|stopword match)\b",
    re.IGNORECASE,
)
LOCAL_PATH = re.compile(r"(?:[A-Za-z]:\\|/Users/|/home/|\\Users\\)", re.IGNORECASE)


@dataclass(frozen=True)
class PowerPointResult:
    path: Path
    slide_count: int
    slide_provenance: list[SlideProvenance]
    evidence_visuals: list[dict[str, Any]]
    validation: dict[str, Any]


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


class PowerPointGenerator:
    def __init__(self, project_root: Path, config: dict[str, Any]):
        self.project_root = project_root.resolve()
        self.config = config
        self.colors = config["colors"]
        self.type = config["typography"]
        self.rtl = False
        self.font = str(self.type["latin_font"])
        self.labels = TEXT["fr"]
        self.provenance: list[SlideProvenance] = []

    @staticmethod
    def _i(value: float | int) -> int:
        return Inches(float(value))

    def _new_presentation(self, language: str) -> Presentation:
        presentation = Presentation()
        presentation.slide_width = self._i(self.config["slide"]["width_inches"])
        presentation.slide_height = self._i(self.config["slide"]["height_inches"])
        self.rtl = language == "ar"
        self.font = str(self.type["arabic_font"] if self.rtl else self.type["latin_font"])
        self.labels = TEXT[language]
        return presentation

    def _shape(
        self,
        slide,
        kind,
        x: float,
        y: float,
        width: float,
        height: float,
        fill: str,
        line: str | None = None,
    ):
        shape = slide.shapes.add_shape(kind, self._i(x), self._i(y), self._i(width), self._i(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(self.colors[fill] if fill in self.colors else fill)
        if line:
            shape.line.color.rgb = _rgb(self.colors[line] if line in self.colors else line)
            shape.line.width = Pt(0.8)
        else:
            shape.line.fill.background()
        shape._element.spPr.get_or_add_effectLst()
        return shape

    @staticmethod
    def _compact_text(value: str, maximum_characters: int) -> str:
        text = re.sub(r"\s+", " ", value).strip(" •▪+-–—:;")
        if len(text) <= maximum_characters:
            return text
        candidate = text[:maximum_characters].rsplit(" ", 1)[0].rstrip(" ,;:-")
        return f"{candidate}…" if candidate else f"{text[:maximum_characters]}…"

    def _text(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        size: float,
        color: str = "body",
        bold: bool = False,
        align: PP_ALIGN | None = None,
        valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
        margin: float = 0.02,
        name: str | None = None,
    ):
        box = slide.shapes.add_textbox(self._i(x), self._i(y), self._i(width), self._i(height))
        if name:
            box.name = name
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = self._i(margin)
        frame.margin_right = self._i(margin)
        frame.margin_top = self._i(margin)
        frame.margin_bottom = self._i(margin)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align or (PP_ALIGN.RIGHT if self.rtl else PP_ALIGN.LEFT)
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        if self.rtl:
            paragraph._p.get_or_add_pPr().set("rtl", "1")
        for run in paragraph.runs:
            run.font.name = self.font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(self.colors[color] if color in self.colors else color)
        return box

    def _bullets(
        self,
        slide,
        bullets: list[str],
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        size: float,
        color: str = "body",
    ):
        box = slide.shapes.add_textbox(self._i(x), self._i(y), self._i(width), self._i(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = self._i(0.05)
        frame.margin_right = self._i(0.03)
        frame.margin_top = self._i(0.02)
        frame.margin_bottom = self._i(0.02)
        for index, text in enumerate(bullets):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = f"▪ {text}"
            paragraph.alignment = PP_ALIGN.RIGHT if self.rtl else PP_ALIGN.LEFT
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(2)
            paragraph.line_spacing = 1.0
            if self.rtl:
                paragraph._p.get_or_add_pPr().set("rtl", "1")
            for run in paragraph.runs:
                run.font.name = self.font
                run.font.size = Pt(size)
                run.font.color.rgb = _rgb(self.colors[color])
        return box

    def _picture_contain(
        self,
        slide,
        path: Path,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        name: str | None = None,
    ):
        with Image.open(path) as image:
            ratio = image.width / image.height
        box_ratio = width / height
        if ratio >= box_ratio:
            rendered_width = width
            rendered_height = width / ratio
        else:
            rendered_height = height
            rendered_width = height * ratio
        picture = slide.shapes.add_picture(
            str(path),
            self._i(x + (width - rendered_width) / 2),
            self._i(y + (height - rendered_height) / 2),
            width=self._i(rendered_width),
            height=self._i(rendered_height),
        )
        if name:
            picture.name = name
        return picture

    def _footer(self, slide, number: int) -> None:
        footer = self.config["footer"]
        logo = (self.project_root / footer["logo_path"]).resolve()
        self._picture_contain(
            slide, logo, footer["logo_x"], footer["logo_y"], footer["logo_width"], 0.28
        )
        self._text(
            slide,
            str(number),
            footer["slide_number_x"],
            footer["slide_number_y"],
            footer["slide_number_width"],
            footer["slide_number_height"],
            size=self.type["slide_number_pt"],
            color="muted",
            bold=True,
            align=PP_ALIGN.CENTER,
            name="Slide number",
        )

    def _title(self, slide, text: str) -> None:
        title = self.config["title"]
        self._text(
            slide,
            text,
            title["x"],
            title["y"],
            title["width"],
            title["height"],
            size=self.type["slide_title_pt"],
            color="heading",
            bold=True,
            name="Slide title",
        )
        self._shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            title["underline_x"],
            title["underline_y"],
            title["underline_width"],
            title["underline_height"],
            "coral",
        )

    def _add_cover(self, presentation: Presentation, request: ReferencePackRequest) -> None:
        cfg = self.config["cover"]
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._shape(slide, MSO_SHAPE.RECTANGLE, cfg["accent_x"], cfg["accent_y"], cfg["accent_width"], cfg["accent_height"], "coral")
        self._shape(slide, MSO_SHAPE.RECTANGLE, cfg["visual_x"], cfg["visual_y"], cfg["visual_width"], cfg["visual_height"], "panel")
        self._shape(slide, MSO_SHAPE.OVAL, cfg["visual_x"] + 0.38, 1.08, 2.25, 2.25, "coral")
        self._shape(slide, MSO_SHAPE.OVAL, cfg["visual_x"] + 0.88, 1.58, 1.25, 1.25, "panel")
        self._text(slide, request.title, cfg["title_x"], cfg["title_y"], cfg["title_width"], cfg["title_height"], size=self.type["cover_title_pt"], color="heading", bold=True, valign=MSO_ANCHOR.MIDDLE, name="Slide title")
        self._text(slide, request.client_name, cfg["client_x"], cfg["client_y"], cfg["client_width"], cfg["client_height"], size=18, color="coral", bold=True)
        if request.subtitle:
            self._text(slide, request.subtitle, cfg["subtitle_x"], cfg["subtitle_y"], cfg["subtitle_width"], cfg["subtitle_height"], size=14, color="body")
        date_text = f"{self.labels['prepared']} {self._date(request.preparation_date, request.language)}"
        self._text(slide, date_text, cfg["date_x"], cfg["date_y"], cfg["date_width"], cfg["date_height"], size=10, color="muted")
        self._footer(slide, 1)
        self.provenance.append(SlideProvenance(slide_number=1, slide_type="cover"))

    @staticmethod
    def _date(value: date, language: str) -> str:
        if language == "en":
            return value.strftime("%Y-%m-%d")
        if language == "ar":
            return value.strftime("%Y/%m/%d")
        return value.strftime("%d/%m/%Y")

    def _add_divider(self, presentation: Presentation) -> None:
        cfg = self.config["section_divider"]
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._text(slide, "2", cfg["number_x"], cfg["number_y"], cfg["number_width"], cfg["number_height"], size=self.type["section_number_pt"], color="coral", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        self._text(slide, self.labels["section"], cfg["title_x"], cfg["title_y"], cfg["title_width"], cfg["title_height"], size=self.type["section_title_pt"], color="heading", bold=False, valign=MSO_ANCHOR.MIDDLE, name="Slide title")
        self._shape(slide, MSO_SHAPE.OVAL, cfg["visual_x"], cfg["visual_y"], cfg["visual_width"], cfg["visual_height"], "coral")
        self._shape(slide, MSO_SHAPE.OVAL, cfg["visual_x"] + 0.53, cfg["visual_y"] + 0.53, cfg["visual_width"] - 1.06, cfg["visual_height"] - 1.06, "white")
        number = len(presentation.slides)
        self._footer(slide, number)
        self.provenance.append(SlideProvenance(slide_number=number, slide_type="section_divider"))

    def _add_summary(self, presentation: Presentation, references: list[PreparedReference]) -> None:
        cfg = self.config["summary"]
        maximum = int(cfg["cards_per_slide"])
        slide_total = (len(references) + maximum - 1) // maximum
        base, remainder = divmod(len(references), slide_total)
        sizes = [base + (1 if index < remainder else 0) for index in range(slide_total)]
        groups: list[list[PreparedReference]] = []
        cursor = 0
        for size in sizes:
            groups.append(references[cursor:cursor + size])
            cursor += size
        for group_index, group in enumerate(groups, start=1):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self._title(slide, self.labels["summary"].format(current=group_index, total=len(groups)))
            bullet_sources = []
            card_height = float(cfg["card_height"]) if len(group) == 3 else 2.30
            card_gap = float(cfg["card_gap"])
            first_y = float(cfg["first_y"]) if len(group) == 3 else 1.40
            for row_index, prepared in enumerate(group):
                y = first_y + row_index * (card_height + card_gap)
                self._shape(slide, MSO_SHAPE.RECTANGLE, cfg["mission_x"], y, cfg["mission_width"], card_height, "coral")
                sequence = references.index(prepared) + 1
                mission = f"{sequence}. {prepared.reference.mission_title}"
                self._text(slide, mission, cfg["mission_x"] + 0.10, y + 0.10, cfg["mission_width"] - 0.20, card_height - 0.20, size=self.type["card_title_pt"], color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                self._shape(slide, MSO_SHAPE.CHEVRON, cfg["mission_x"] + cfg["mission_width"] - 0.05, y + (card_height - 0.58) / 2, 0.27, 0.58, "coral")
                self._shape(slide, MSO_SHAPE.RECTANGLE, cfg["content_x"], y, cfg["content_width"], card_height, "white", "border")
                bullets = [
                    self._compact_text(item.text, int(cfg["maximum_bullet_characters"]))
                    for item in prepared.summary_bullets[: int(cfg["maximum_bullets"])]
                ]
                self._bullets(slide, bullets, cfg["content_x"] + 0.15, y + 0.12, cfg["content_width"] - 0.30, card_height - 0.24, size=self.type["body_min_pt"])
                self._shape(slide, MSO_SHAPE.RECTANGLE, cfg["client_x"], y, cfg["client_width"], card_height, "white", "border")
                client_meta = "\n".join(value for value in (prepared.reference.client, prepared.reference.country, prepared.reference.period) if value)
                self._text(slide, client_meta, cfg["client_x"] + 0.10, y + 0.12, cfg["client_width"] - 0.20, card_height - 0.24, size=self.type["metadata_pt"], color="body", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
                bullet_sources.extend(prepared.summary_bullets[: int(cfg["maximum_bullets"])])
            number = len(presentation.slides)
            self._footer(slide, number)
            self.provenance.append(SlideProvenance(slide_number=number, slide_type="reference_summary", reference_ids=[item.reference.reference_id for item in group], bullet_sources=bullet_sources))

    def _metadata_block(self, slide, prepared: PreparedReference, x: float, y: float, width: float, height: float) -> None:
        reference = prepared.reference
        values = [
            (self.labels["client"], reference.client),
            (self.labels["country"], reference.country),
            (self.labels["period"], reference.period),
            (self.labels["sector"], reference.sector),
            (self.labels["offering"], reference.offering),
        ]
        self._shape(slide, MSO_SHAPE.RECTANGLE, x, y, width, height, "panel", "border")
        cursor = y + 0.18
        for label, value in values:
            if not value:
                continue
            self._text(slide, label.upper(), x + 0.18, cursor, width - 0.36, 0.22, size=8, color="coral", bold=True)
            self._text(slide, value, x + 0.18, cursor + 0.23, width - 0.36, 0.42, size=11, color="heading", bold=True)
            cursor += 0.84

    def _add_details(self, presentation: Presentation, references: list[PreparedReference]) -> None:
        cfg = self.config["detail"]
        for index, prepared in enumerate(references, start=1):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self._title(slide, self.labels["detail"].format(current=index, total=len(references)))
            self._metadata_block(slide, prepared, cfg["left_x"], cfg["content_y"], cfg["left_width"], cfg["content_height"])
            right_x = float(cfg["right_x"])
            right_width = float(cfg["right_width"])
            self._text(slide, prepared.reference.mission_title, right_x, cfg["content_y"], right_width, 0.72, size=17, color="heading", bold=True)
            self._text(slide, self.labels["services"].upper(), right_x, cfg["content_y"] + 0.88, right_width, 0.24, size=9, color="coral", bold=True)
            service_text = [item.text for item in prepared.service_items[: int(cfg["maximum_service_bullets"])]]
            self._bullets(slide, service_text, right_x, cfg["content_y"] + 1.16, right_width, 1.74, size=self.type["body_pt"])
            technology_text = list(dict.fromkeys([*prepared.reference.technologies, *prepared.reference.capabilities]))
            if technology_text:
                self._text(slide, self.labels["technologies"].upper(), right_x, cfg["content_y"] + 3.00, right_width, 0.24, size=9, color="coral", bold=True)
                self._text(slide, " · ".join(technology_text), right_x, cfg["content_y"] + 3.28, right_width, 0.42, size=10, color="body")
            self._shape(slide, MSO_SHAPE.RECTANGLE, right_x, cfg["content_y"] + 3.82, right_width, 1.14, "panel_alt", "border")
            self._text(slide, self.labels["why"].upper(), right_x + 0.15, cfg["content_y"] + 3.96, right_width - 0.30, 0.22, size=9, color="coral", bold=True)
            self._bullets(slide, [item.text for item in prepared.why_selected], right_x + 0.15, cfg["content_y"] + 4.22, right_width - 0.30, 0.58, size=9)
            citations = list(dict.fromkeys(item.citation_label for item in prepared.reference.evidence[:2]))
            self._text(slide, f"{self.labels['source']}: " + " · ".join(citations), right_x, cfg["content_y"] + 5.13, right_width, 0.35, size=8, color="muted")
            number = len(presentation.slides)
            self._footer(slide, number)
            sources = [*prepared.service_items[: int(cfg["maximum_service_bullets"])], *prepared.why_selected]
            self.provenance.append(SlideProvenance(slide_number=number, slide_type="reference_detail", reference_ids=[prepared.reference.reference_id], bullet_sources=sources, evidence_chunk_ids=[item.chunk_id for item in prepared.reference.evidence[:2]]))

    def _evidence_fallback(
        self,
        slide,
        evidence,
        visual: EvidenceVisual,
    ) -> None:
        cfg = self.config["evidence"]
        x, y = float(cfg["image_x"]), float(cfg["image_y"])
        width, height = float(cfg["image_width"]), float(cfg["image_height"])
        self._shape(slide, MSO_SHAPE.RECTANGLE, x, y, width, height, "panel_alt", "coral")
        self._shape(slide, MSO_SHAPE.RECTANGLE, x, y, width, 0.58, "coral")
        self._text(
            slide,
            self.labels["fallback"],
            x + 0.28,
            y + 0.13,
            width - 0.56,
            0.30,
            size=13,
            color="white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        excerpt = prepare_evidence_excerpt(evidence.display_text, int(cfg["maximum_characters"]))
        self._text(
            slide,
            excerpt,
            x + 0.60,
            y + 0.90,
            width - 1.20,
            height - 1.58,
            size=15,
            color="body",
            valign=MSO_ANCHOR.MIDDLE,
            name="Evidence text fallback",
        )
        self._text(
            slide,
            visual.fallback_reason or "source page rendering unavailable",
            x + 0.60,
            y + height - 0.52,
            width - 1.20,
            0.24,
            size=8,
            color="muted",
            align=PP_ALIGN.CENTER,
        )

    def _add_evidence(
        self,
        presentation: Presentation,
        references: list[PreparedReference],
        visuals: dict[str, EvidenceVisual],
    ) -> None:
        cfg = self.config["evidence"]
        total = len(references)
        for current, prepared in enumerate(references, start=1):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self._title(slide, self.labels["evidence"].format(current=current, total=total))
            evidence = prepared.evidence_items[0]
            visual = visuals[prepared.reference.reference_id]
            self._text(
                slide,
                prepared.reference.mission_title,
                0.78,
                cfg["content_y"],
                11.77,
                cfg["reference_title_height"],
                size=int(self.type.get("evidence_title_pt", 15)),
                color="heading",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                name="Evidence reference title",
            )
            subtitle = " · ".join(
                value for value in (prepared.reference.client, prepared.reference.country) if value
            )
            self._text(
                slide,
                subtitle,
                0.78,
                cfg["subtitle_y"],
                11.77,
                cfg["subtitle_height"],
                size=9,
                color="coral",
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            if visual.is_rendered:
                source_ratio = visual.rendered_pixel_width / visual.rendered_pixel_height
                box_ratio = float(cfg["image_width"]) / float(cfg["image_height"])
                if source_ratio >= box_ratio:
                    frame_width = float(cfg["image_width"])
                    frame_height = frame_width / source_ratio
                else:
                    frame_height = float(cfg["image_height"])
                    frame_width = frame_height * source_ratio
                frame_x = float(cfg["image_x"]) + (float(cfg["image_width"]) - frame_width) / 2
                frame_y = float(cfg["image_y"]) + (float(cfg["image_height"]) - frame_height) / 2
                self._shape(
                    slide,
                    MSO_SHAPE.RECTANGLE,
                    frame_x - 0.04,
                    frame_y - 0.04,
                    frame_width + 0.08,
                    frame_height + 0.08,
                    "white",
                    "coral",
                )
                picture = self._picture_contain(
                    slide,
                    visual.image_path,
                    cfg["image_x"],
                    cfg["image_y"],
                    cfg["image_width"],
                    cfg["image_height"],
                    name="Evidence source image",
                )
                expected_ratio = visual.rendered_pixel_width / visual.rendered_pixel_height
                placed_ratio = picture.width / picture.height
                record = visual.manifest_record()
                record["placed_aspect_ratio"] = round(placed_ratio, 6)
                record["aspect_ratio_preserved"] = abs(placed_ratio - expected_ratio) <= 0.01
            else:
                self._evidence_fallback(slide, evidence, visual)
                record = visual.manifest_record()
            citation = (
                f"{self.labels['source']}: {evidence.source_file_name} · "
                f"{self.labels['page']} {evidence.source_page}"
            )
            self._text(
                slide,
                citation,
                0.78,
                cfg["citation_y"],
                11.77,
                cfg["citation_height"],
                size=8,
                color="muted",
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
            )
            number = len(presentation.slides)
            self._footer(slide, number)
            self.provenance.append(
                SlideProvenance(
                    slide_number=number,
                    slide_type="evidence_annex",
                    reference_ids=[prepared.reference.reference_id],
                    evidence_chunk_ids=[evidence.chunk_id],
                    evidence_visuals=[record],
                )
            )

    def generate(self, output_path: Path, request: ReferencePackRequest, references: list[PreparedReference]) -> PowerPointResult:
        self.provenance = []
        visuals: dict[str, EvidenceVisual] = {}
        if request.include_evidence_annex:
            renderer = EvidenceRenderer(self.project_root, self.config)
            for index, prepared in enumerate(references, start=1):
                evidence = prepared.evidence_items[0]
                visuals[prepared.reference.reference_id] = renderer.render(
                    prepared.reference.reference_id,
                    evidence,
                    output_path.parent
                    / "evidence_images"
                    / f"reference-{index:03d}-page-{evidence.source_page}.png",
                )
        presentation = self._new_presentation(request.language)
        self._add_cover(presentation, request)
        self._add_divider(presentation)
        if request.include_summary:
            self._add_summary(presentation, references)
        if request.include_reference_details:
            self._add_details(presentation, references)
        if request.include_evidence_annex:
            self._add_evidence(presentation, references, visuals)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        validation = self.validate(output_path, len(references))
        evidence_visuals = [
            visual
            for provenance in self.provenance
            for visual in provenance.evidence_visuals
        ]
        return PowerPointResult(
            output_path,
            len(presentation.slides),
            self.provenance,
            evidence_visuals,
            validation,
        )

    def validate(self, path: Path, selected_count: int) -> dict[str, Any]:
        presentation = Presentation(path)
        width, height = presentation.slide_width, presentation.slide_height
        errors: list[str] = []
        titles = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            has_title = False
            has_number = False
            for shape in slide.shapes:
                if shape.left < 0 or shape.top < 0 or shape.left + shape.width > width or shape.top + shape.height > height:
                    errors.append(f"slide {slide_index}: shape outside slide bounds")
                if getattr(shape, "name", "") == "Slide title":
                    has_title = True
                if getattr(shape, "name", "") == "Slide number":
                    has_number = True
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text
                if PROHIBITED_TEXT.search(text):
                    errors.append(f"slide {slide_index}: internal retrieval score text detected")
                if LOCAL_PATH.search(text):
                    errors.append(f"slide {slide_index}: local filesystem path detected")
            if slide_index > 1 and not has_title:
                errors.append(f"slide {slide_index}: missing title")
            else:
                titles += int(has_title)
            if not has_number:
                errors.append(f"slide {slide_index}: missing slide number")
        if len({item.reference_ids[0] for item in self.provenance if item.slide_type == "reference_detail" and item.reference_ids}) not in {0, selected_count}:
            errors.append("detailed reference slides are duplicated or missing")
        for evidence_slide in (
            item for item in self.provenance if item.slide_type == "evidence_annex"
        ):
            if len(evidence_slide.reference_ids) != 1 or len(evidence_slide.evidence_visuals) != 1:
                errors.append(
                    f"slide {evidence_slide.slide_number}: evidence ownership is ambiguous"
                )
                continue
            visual = evidence_slide.evidence_visuals[0]
            if not visual.get("rendered_source_image") and not visual.get("fallback_reason"):
                errors.append(
                    f"slide {evidence_slide.slide_number}: no source image or documented fallback"
                )
            if visual.get("rendered_source_image") and not visual.get("aspect_ratio_preserved"):
                errors.append(
                    f"slide {evidence_slide.slide_number}: source image aspect ratio changed"
                )
        if errors:
            raise RuntimeError("PowerPoint validation failed: " + "; ".join(errors))
        return {
            "status": "PASS",
            "slide_count": len(presentation.slides),
            "selected_reference_count": selected_count,
            "titles_present": titles,
            "bounds_check": "PASS",
            "internal_score_scan": "PASS",
            "local_path_scan": "PASS",
            "slide_numbering": "PASS",
            "evidence_visuals": "PASS",
            "evidence_aspect_ratio": "PASS",
        }
