from __future__ import annotations

import json
import math
import shutil
from pathlib import Path

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation

from app.api.settings import load_config
from reference_pack.schemas import ReferencePackRequest
from reference_pack.service import ReferencePackService


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "generated/reference_packs/design_correction_demo"
AUDIT = ROOT / "audit/reference_pack/design_correction_demo"
REFERENCE_IDS = [
    "6f1f43d3bfa0132fe5d5a08c4a17868e0521f514a0336d065a227518e2a3e6dd",
    "e2b45d3d3e5a0684286f6b7348f91937f66d6308fd13edb60a9665b4ef4b0a93",
    "a08cc6dd4732fbda1346b545cd98fd076524375a5ccf6803fd1afe37ebe1271b",
    "2650c62de0abf30996b78316e109aab2ebd98b5f2e8c65db10943f1785fcb982",
]


def _render_page(pdf: Path, index: int, output: Path, scale: float = 1.6) -> None:
    with fitz.open(pdf) as document:
        pixmap = document[index].get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        pixmap.save(output)


def _preserve_old_evidence() -> str | None:
    inventory = ROOT / "audit/reference_pack/demo_generations.json"
    if not inventory.is_file():
        return None
    payload = json.loads(inventory.read_text(encoding="utf-8"))
    old_case = next((case for case in payload.get("outputs", []) if case.get("case") == "ar-four"), None)
    if not old_case:
        return None
    directory = Path(old_case["directory"])
    manifest_path = directory / "generation_manifest.json"
    pdf_path = directory / "reference_pack.pdf"
    if not manifest_path.is_file() or not pdf_path.is_file():
        return None
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    evidence = next(
        (slide for slide in manifest.get("slide_provenance", []) if slide.get("slide_type") == "evidence_annex"),
        None,
    )
    if not evidence:
        return None
    output = AUDIT / "old_evidence_slide.png"
    _render_page(pdf_path, int(evidence["slide_number"]) - 1, output)
    return output.relative_to(ROOT).as_posix()


def _contact_sheet(slides: list[Path], output: Path) -> None:
    columns = 3
    thumb_width, thumb_height, label_height, gap = 480, 270, 24, 16
    rows = math.ceil(len(slides) / columns)
    sheet = Image.new(
        "RGB",
        (columns * (thumb_width + gap) + gap, rows * (thumb_height + label_height + gap) + gap),
        "#ececef",
    )
    draw = ImageDraw.Draw(sheet)
    for index, path in enumerate(slides):
        with Image.open(path) as opened:
            image = opened.convert("RGB")
        image.thumbnail((thumb_width, thumb_height), Image.Resampling.LANCZOS)
        column, row = index % columns, index // columns
        x = gap + column * (thumb_width + gap)
        y = gap + row * (thumb_height + label_height + gap)
        draw.text((x, y + 4), f"Slide {index + 1}", fill="#161616")
        sheet.paste(image, (x, y + label_height))
    sheet.save(output, optimize=True)


def _validate(pptx: Path, pdf: Path, manifest: dict) -> dict:
    presentation = Presentation(pptx)
    with fitz.open(pdf) as document:
        pdf_pages = document.page_count
    errors: list[str] = []
    evidence = [
        slide for slide in manifest["slide_provenance"] if slide["slide_type"] == "evidence_annex"
    ]
    summaries = [
        slide for slide in manifest["slide_provenance"] if slide["slide_type"] == "reference_summary"
    ]
    if len(evidence) != len(REFERENCE_IDS):
        errors.append("one evidence slide per selected reference is required")
    if any(len(slide["reference_ids"]) > 3 for slide in summaries):
        errors.append("a summary slide contains more than three references")
    for slide_record in evidence:
        number = int(slide_record["slide_number"])
        slide = presentation.slides[number - 1]
        visuals = slide_record.get("evidence_visuals", [])
        if len(visuals) != 1:
            errors.append(f"slide {number}: missing evidence visual record")
            continue
        visual = visuals[0]
        text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
        if visual.get("rendered_source_image"):
            if not any(shape.name == "Evidence source image" for shape in slide.shapes):
                errors.append(f"slide {number}: source image missing")
            if not visual.get("aspect_ratio_preserved"):
                errors.append(f"slide {number}: image aspect ratio changed")
        elif not visual.get("fallback_reason"):
            errors.append(f"slide {number}: undocumented fallback")
        if visual["source_file_name"] not in text or f"page {visual['source_page']}" not in text:
            errors.append(f"slide {number}: source filename/page citation missing")
    if len(presentation.slides) != pdf_pages:
        errors.append("PPTX/PDF slide count mismatch")
    if errors:
        raise RuntimeError("; ".join(errors))
    return {
        "status": "PASS",
        "selected_reference_count": len(REFERENCE_IDS),
        "slide_count": len(presentation.slides),
        "evidence_slide_count": len(evidence),
        "rendered_source_page_count": sum(
            bool(slide["evidence_visuals"][0].get("rendered_source_image")) for slide in evidence
        ),
        "documented_fallback_count": sum(
            bool(slide["evidence_visuals"][0].get("fallback_reason")) for slide in evidence
        ),
        "summary_maximum_references": max(len(slide["reference_ids"]) for slide in summaries),
        "aspect_ratio": "PASS",
        "source_citations": "PASS",
        "pptx_pdf_page_parity": "PASS",
    }


def main() -> int:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    AUDIT.mkdir(parents=True, exist_ok=True)
    old_evidence = _preserve_old_evidence()

    service = ReferencePackService(ROOT, load_config())
    artifact = service.generate(
        ReferencePackRequest(
            title="Références Devoteam pertinentes",
            client_name="Démonstration design corrigé",
            subtitle="Quatre références approuvées — source-page evidence",
            language="fr",
            reference_ids=REFERENCE_IDS,
            output_formats=["pptx", "pdf"],
        )
    )
    generated = Path(artifact.directory)
    for name in ("reference_pack.pptx", "reference_pack.pdf", "generation_manifest.json"):
        target = OUTPUT / name
        if target.is_file():
            target.unlink()
        shutil.copy2(generated / name, target)

    manifest = json.loads((OUTPUT / "generation_manifest.json").read_text(encoding="utf-8"))
    validation = _validate(OUTPUT / "reference_pack.pptx", OUTPUT / "reference_pack.pdf", manifest)

    for path in AUDIT.glob("slide-*.png"):
        path.unlink()
    slides: list[Path] = []
    with fitz.open(OUTPUT / "reference_pack.pdf") as document:
        for index, page in enumerate(document):
            path = AUDIT / f"slide-{index + 1:03d}.png"
            page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False).save(path)
            slides.append(path)
    _contact_sheet(slides, AUDIT / "contact-sheet.png")
    first_evidence = next(
        slide for slide in manifest["slide_provenance"] if slide["slide_type"] == "evidence_annex"
    )
    shutil.copy2(slides[int(first_evidence["slide_number"]) - 1], AUDIT / "corrected_evidence_slide.png")
    validation["old_evidence_slide"] = old_evidence
    validation["corrected_evidence_slide"] = (
        AUDIT / "corrected_evidence_slide.png"
    ).relative_to(ROOT).as_posix()
    (AUDIT / "validation.json").write_text(
        json.dumps(validation, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    print(json.dumps(validation, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
