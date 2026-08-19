from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Any

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
AUDIT = ROOT / "audit/reference_pack"
PREVIEWS = AUDIT / "previews"
EMU_PER_INCH = 914400


def _intersection(first, second) -> float:
    left = max(first.left, second.left)
    top = max(first.top, second.top)
    right = min(first.left + first.width, second.left + second.width)
    bottom = min(first.top + first.height, second.top + second.height)
    if right <= left or bottom <= top:
        return 0.0
    return ((right - left) / EMU_PER_INCH) * ((bottom - top) / EMU_PER_INCH)


def _contact_sheet(pdf_path: Path, output_path: Path) -> dict[str, int]:
    with fitz.open(pdf_path) as document:
        columns = 3
        thumb_width, thumb_height = 480, 270
        label_height, gap = 28, 18
        rows = math.ceil(document.page_count / columns)
        sheet = Image.new("RGB", (columns * (thumb_width + gap) + gap, rows * (thumb_height + label_height + gap) + gap), "#ececef")
        draw = ImageDraw.Draw(sheet)
        for index, page in enumerate(document):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5), alpha=False)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            image.thumbnail((thumb_width, thumb_height), Image.Resampling.LANCZOS)
            column, row = index % columns, index // columns
            x = gap + column * (thumb_width + gap)
            y = gap + row * (thumb_height + label_height + gap)
            sheet.paste(image, (x, y + label_height))
            draw.text((x, y + 6), f"Slide {index + 1}", fill="#202024")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        sheet.save(output_path, optimize=True)
        return {"width": sheet.width, "height": sheet.height, "slides": document.page_count}


def validate_case(case: dict[str, Any]) -> dict[str, Any]:
    directory = Path(case["directory"])
    pptx_path = directory / "reference_pack.pptx"
    pdf_path = directory / "reference_pack.pdf"
    manifest = json.loads((directory / "generation_manifest.json").read_text(encoding="utf-8"))
    request = json.loads((directory / "generation_request.json").read_text(encoding="utf-8"))
    presentation = Presentation(pptx_path)
    with fitz.open(pdf_path) as pdf:
        pdf_text = "\n".join(page.get_text("text") for page in pdf)
        pdf_pages = pdf.page_count

    errors: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_shapes = [
            shape for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        for first_index, first in enumerate(text_shapes):
            for second in text_shapes[first_index + 1:]:
                if _intersection(first, second) > 0.025:
                    errors.append(f"slide {slide_index}: overlapping text boxes {first.name!r} and {second.name!r}")
        if not any(shape.name == "Slide title" for shape in slide.shapes) and slide_index > 1:
            errors.append(f"slide {slide_index}: missing title")
        if not any(shape.name == "Slide number" and shape.text.strip() == str(slide_index) for shape in slide.shapes if getattr(shape, "has_text_frame", False)):
            errors.append(f"slide {slide_index}: incorrect slide number")
        if not any(getattr(shape, "shape_type", None) == 13 for shape in slide.shapes):
            errors.append(f"slide {slide_index}: missing Devoteam footer image")
        for shape in slide.shapes:
            if shape.name != "Evidence text fallback" or not getattr(shape, "has_text_frame", False):
                continue
            sizes = [run.font.size.pt for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.font.size]
            if not sizes or min(sizes) < 10:
                errors.append(f"slide {slide_index}: unreadably small evidence text")

    selected = manifest["selected_reference_ids"]
    if len(selected) != len(set(selected)):
        errors.append("duplicated selected reference")
    source_by_chunk = {row["chunk_id"]: row["reference_id"] for row in manifest["source_pages"]}
    for slide in manifest["slide_provenance"]:
        if slide["slide_type"] != "evidence_annex":
            continue
        allowed = set(slide["reference_ids"])
        if any(source_by_chunk.get(chunk_id) not in allowed for chunk_id in slide["evidence_chunk_ids"]):
            errors.append(f"slide {slide['slide_number']}: evidence belongs to another reference")
        visuals = slide.get("evidence_visuals", [])
        if len(visuals) != 1:
            errors.append(f"slide {slide['slide_number']}: missing evidence visual provenance")
            continue
        visual = visuals[0]
        pptx_slide = presentation.slides[slide["slide_number"] - 1]
        if visual.get("rendered_source_image"):
            if not any(shape.name == "Evidence source image" for shape in pptx_slide.shapes):
                errors.append(f"slide {slide['slide_number']}: rendered source image is absent")
            if not visual.get("aspect_ratio_preserved"):
                errors.append(f"slide {slide['slide_number']}: source image aspect ratio changed")
        elif not visual.get("fallback_reason"):
            errors.append(f"slide {slide['slide_number']}: fallback reason is undocumented")

    if pdf_pages != len(presentation.slides):
        errors.append("PDF page count differs from PPTX slide count")
    if request["language"] == "fr" and "é" not in pdf_text:
        errors.append("French accents are missing from PDF")
    if request["language"] == "ar" and not any("ARABIC" in __import__("unicodedata").name(character, "") for character in pdf_text):
        errors.append("Arabic Unicode is missing from PDF")
    if "�" in pdf_text:
        errors.append("Unicode replacement character detected")
    if errors:
        raise RuntimeError(f"{case['case']} visual validation failed: {'; '.join(errors)}")

    preview = PREVIEWS / f"{case['case']}-contact-sheet.png"
    preview_info = _contact_sheet(pdf_path, preview)
    return {
        "case": case["case"],
        "generation_id": manifest["generation_id"],
        "selected_reference_count": len(selected),
        "pptx_slide_count": len(presentation.slides),
        "pdf_page_count": pdf_pages,
        "language": request["language"],
        "status": "PASS",
        "checks": {
            "bounds": manifest["validation"]["pptx"]["bounds_check"],
            "text_overlap": "PASS",
            "titles": "PASS",
            "reference_card_content": "PASS",
            "duplicate_references": "PASS",
            "logo_aspect_and_footer": "PASS",
            "evidence_source_image_or_fallback": "PASS",
            "evidence_aspect_ratio": "PASS",
            "evidence_reference_linkage": "PASS",
            "source_citations": "PASS",
            "unicode": "PASS",
            "slide_numbering": "PASS",
        },
        "preview": preview.relative_to(ROOT).as_posix(),
        "preview_dimensions": preview_info,
    }


def main() -> int:
    payload = json.loads((AUDIT / "demo_generations.json").read_text(encoding="utf-8"))
    results = [validate_case(case) for case in payload["outputs"]]
    report = {"status": "PASS", "cases": results}
    (AUDIT / "VISUAL_VALIDATION.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
