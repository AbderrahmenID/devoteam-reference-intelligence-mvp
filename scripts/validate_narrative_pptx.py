from __future__ import annotations

import argparse
import hashlib
import io
import json
import re
import zipfile
from pathlib import Path

import fitz
from PIL import Image
from pptx import Presentation


SAMPLE_MEDIA_SHA1 = {
    "51939cfa864b25a3fd815b8d0494953f47e58c14",
    "964cfe8064569051db76509de55ea7d63dc4729c",
    "90edb490e65f5ad7e96f8c854122b31c15df3c00",
}
SAMPLE_TEXT = re.compile(
    r"mewa|mhrsd|water sustainability|accelerator|strategic domain",
    re.IGNORECASE,
)
INTERNAL_TEXT = re.compile(
    r"(?:[A-Za-z]:\\|/Users/|/home/|\\Users\\|chunk[_ -]?id|retrieval score|embedding[_ -]?id)",
    re.IGNORECASE,
)


def walk(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from walk(shape.shapes)


def validate(path: Path, pdf_path: Path | None = None) -> dict:
    presentation = Presentation(path)
    width, height = presentation.slide_width, presentation.slide_height
    failures: list[dict] = []
    slides: list[dict] = []
    deck_is_compact = any(
        shape.name.startswith("C.")
        for slide in presentation.slides
        for shape in walk(slide.shapes)
    )
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text = "\n".join(shape.text for shape in walk(slide.shapes) if hasattr(shape, "text"))
        names = {shape.name for shape in walk(slide.shapes)}
        is_compact = deck_is_compact
        is_evidence = bool({"D.EVIDENCE_SOURCE_IMAGE", "Evidence source image"} & names)
        if is_compact:
            required = {"Slide title", "Slide number"}
            if slide_number == 1:
                required |= {"C.SECTION_NUMBER"}
            elif not is_evidence:
                if not any(name.startswith("C.MISSION_") for name in names):
                    failures.append({"slide": slide_number, "reason": "MISSING_COMPACT_REFERENCE_CARD"})
        else:
            required = {"D.DEVOTEAM_LOGO", "D.PAGE_NUMBER"}
            if "D.MISSION_TITLE" in names:
                required |= {
                    "D.MISSION_TITLE",
                    "D.CLIENT",
                    "D.COUNTRY",
                    "D.OFFERING",
                    "D.SECTOR",
                    "D.PERIOD",
                    "D.CHALLENGE",
                    "D.REALISATIONS",
                    "D.BENEFITS",
                }
            elif slide_number == 1:
                required |= {"D.SECTION_TITLE"}
            elif is_evidence:
                required |= {"D.EVIDENCE_TITLE", "D.EVIDENCE_METADATA", "D.EVIDENCE_SOURCE_LABEL"}
        missing = sorted(required - names)
        if missing:
            failures.append({"slide": slide_number, "reason": "MISSING_EDITABLE_SHAPE", "shapes": missing})
        if SAMPLE_TEXT.search(text):
            failures.append({"slide": slide_number, "reason": "SAMPLE_TEXT_LEAK"})
        if INTERNAL_TEXT.search(text):
            failures.append({"slide": slide_number, "reason": "INTERNAL_TEXT_LEAK"})
        out_of_bounds = []
        fonts_below_minimum = []
        bullet_failures = []
        for shape in walk(slide.shapes):
            if (
                shape.name != "Orange source bleed"
                and (
                    shape.left < 0
                    or shape.top < 0
                    or shape.left + shape.width > width
                    or shape.top + shape.height > height
                )
            ):
                out_of_bounds.append(shape.name)
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size.pt < 8:
                            fonts_below_minimum.append({"shape": shape.name, "size_pt": run.font.size.pt})
                if shape.name in {"D.CHALLENGE", "D.REALISATIONS", "D.BENEFITS"}:
                    for paragraph in shape.text_frame.paragraphs[1:]:
                        if paragraph.text.strip() and not paragraph._p.xpath("./a:pPr/a:buChar"):
                            bullet_failures.append(shape.name)
        if out_of_bounds:
            failures.append({"slide": slide_number, "reason": "OUT_OF_BOUNDS", "shapes": out_of_bounds})
        if fonts_below_minimum:
            failures.append({"slide": slide_number, "reason": "FONT_BELOW_MINIMUM", "runs": fonts_below_minimum})
        if bullet_failures:
            failures.append({"slide": slide_number, "reason": "NON_NATIVE_BULLET", "shapes": sorted(set(bullet_failures))})
        evidence_aspect_ratio_pass = True
        if is_evidence:
            pictures = [
                shape
                for shape in walk(slide.shapes)
                if shape.name in {"D.EVIDENCE_SOURCE_IMAGE", "Evidence source image"}
            ]
            if len(pictures) != 1:
                failures.append({"slide": slide_number, "reason": "EVIDENCE_IMAGE_COUNT", "count": len(pictures)})
                evidence_aspect_ratio_pass = False
            else:
                picture = pictures[0]
                with Image.open(io.BytesIO(picture.image.blob)) as image:
                    source_ratio = image.width / image.height
                placed_ratio = picture.width / picture.height
                evidence_aspect_ratio_pass = abs(source_ratio - placed_ratio) <= 0.01
                if not evidence_aspect_ratio_pass:
                    failures.append({"slide": slide_number, "reason": "EVIDENCE_ASPECT_RATIO_CHANGED"})
        slides.append(
            {
                "slide_number": slide_number,
                "editable_text_shape_count": sum(1 for shape in walk(slide.shapes) if getattr(shape, "has_text_frame", False)),
                "out_of_bounds_shape_count": len(out_of_bounds),
                "minimum_font_pass": not fonts_below_minimum,
                "native_bullet_pass": not bullet_failures,
                "sample_text_pass": not SAMPLE_TEXT.search(text),
                "internal_text_pass": not INTERNAL_TEXT.search(text),
                "evidence_slide": is_evidence,
                "evidence_aspect_ratio_pass": evidence_aspect_ratio_pass,
            }
        )
    with zipfile.ZipFile(path) as archive:
        media_hashes = {
            hashlib.sha1(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith("ppt/media/")
        }
        slide_xml = [name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)]
        empty_placeholders = []
        package_text_leaks = []
        for name in archive.namelist():
            if name.endswith((".xml", ".rels", ".txt")):
                decoded = archive.read(name).decode("utf-8", errors="ignore")
                if SAMPLE_TEXT.search(decoded):
                    package_text_leaks.append(name)
        for name in slide_xml:
            xml = archive.read(name).decode("utf-8")
            for placeholder in re.findall(r"<p:sp>.*?</p:sp>", xml, flags=re.DOTALL):
                if "<p:ph" in placeholder and not any(value.strip() for value in re.findall(r"<a:t>(.*?)</a:t>", placeholder, flags=re.DOTALL)):
                    empty_placeholders.append(name)
        if media_hashes & SAMPLE_MEDIA_SHA1:
            failures.append({"reason": "SAMPLE_MEDIA_LEAK", "sha1": sorted(media_hashes & SAMPLE_MEDIA_SHA1)})
        if empty_placeholders:
            failures.append({"reason": "EMPTY_STRUCTURAL_PLACEHOLDER", "slides": sorted(set(empty_placeholders))})
        if package_text_leaks:
            failures.append({"reason": "SAMPLE_PACKAGE_TEXT_LEAK", "parts": sorted(package_text_leaks)})
    pdf_validation: dict = {"status": "NOT_REQUESTED"}
    if pdf_path is not None:
        if not pdf_path.is_file():
            failures.append({"reason": "PDF_NOT_FOUND"})
            pdf_validation = {"status": "FAIL", "reason": "PDF_NOT_FOUND"}
        else:
            with fitz.open(pdf_path) as document:
                blank_pages = []
                for index, page in enumerate(document, start=1):
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(0.25, 0.25), alpha=False)
                    if not any(value < 245 for value in pixmap.samples):
                        blank_pages.append(index)
                pdf_validation = {
                    "status": "PASS" if document.page_count == len(presentation.slides) and not blank_pages else "FAIL",
                    "page_count": document.page_count,
                    "blank_pages": blank_pages,
                }
            if pdf_validation["status"] != "PASS":
                failures.append({"reason": "PDF_PPTX_MISMATCH", **pdf_validation})
    return {
        "status": "PASS" if not failures else "FAIL",
        "pptx": path.name,
        "slide_count": len(presentation.slides),
        "slide_size_inches": [round(width / 914400, 3), round(height / 914400, 3)],
        "slides": slides,
        "sample_media_pass": not bool(media_hashes & SAMPLE_MEDIA_SHA1),
        "empty_placeholder_pass": not empty_placeholders,
        "package_text_pass": not package_text_leaks,
        "pdf_validation": pdf_validation,
        "failures": failures,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--pdf", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()
    report = validate(args.pptx.resolve(), args.pdf.resolve() if args.pdf else None)
    rendered = json.dumps(report, ensure_ascii=False, indent=2) + "\n"
    if args.output:
        args.output.write_text(rendered, encoding="utf-8")
    print(rendered, end="")
    return 0 if report["status"] == "PASS" else 1


if __name__ == "__main__":
    raise SystemExit(main())
