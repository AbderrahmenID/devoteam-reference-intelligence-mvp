from __future__ import annotations

import hashlib
import json
from pathlib import Path

import fitz
import pytest
from PIL import Image
from pydantic import ValidationError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from reference_narrative.compact_pptx_generator import (
    ORANGE_COMPACT_SOURCE_SLIDES,
    ORANGE_DIVIDER_SOURCE_SLIDE,
    ORANGE_EVIDENCE_SOURCE_SLIDES,
    ORANGE_SOURCE_FIRST_SLIDE,
    ORANGE_SOURCE_LAST_SLIDE,
    OrangeBankCompactNarrativePptxGenerator,
)
from reference_narrative.evidence_annex import NarrativeEvidenceSelection
from reference_narrative.presentation_schemas import (
    SUPPORTED_TEMPLATE_IDS,
    NarrativePresentationRequest,
)
from reference_narrative.template_mapper import PptxContentOverflowError
from reference_pack.evidence_renderer import EvidenceVisual
from reference_pack.schemas import TrustedEvidence, TrustedReference
from reference_pack.validation import sha256_file

from test_narrative_pptx import REF_A, _ids, _request, _review, _service, _service_request, _shape, _text


def _selections(ids: list[str], image_path: Path) -> list[NarrativeEvidenceSelection]:
    with Image.open(image_path) as image:
        width, height = image.size
    selections: list[NarrativeEvidenceSelection] = []
    for index, reference_id in enumerate(ids, start=1):
        evidence = TrustedEvidence(
            chunk_id=f"chunk-{index}",
            document_id=f"document-{index}",
            source_file_name=f"approved-source-{index}.pdf",
            source_sha256=sha256_file(image_path),
            source_page=1,
            citation_label=f"Approved source {index} - page 1",
            citation_uri="https://example.test/source",
            language="fr",
            display_text=f"Approved evidence {index}",
        )
        reference = TrustedReference(
            reference_id=reference_id,
            row_number=index,
            mission_title=f"Trusted mission {index}",
            client=f"Trusted client {index}",
            country=f"Trusted country {index}",
            period=f"20{20 + index}",
            sector=f"Trusted sector {index}",
            offering=f"Trusted offering {index}",
            business_unit="Cloud",
            description="Approved reference",
            services_delivered=[],
            technologies=[],
            capabilities=[],
            evidence=[evidence],
        )
        visual = EvidenceVisual(
            reference_id=reference_id,
            chunk_id=evidence.chunk_id,
            document_id=evidence.document_id,
            source_file_name=evidence.source_file_name,
            source_page=1,
            source_sha256=evidence.source_sha256,
            image_path=image_path,
            rendering_method="source_image_render",
            fallback_reason=None,
            original_pixel_width=width,
            original_pixel_height=height,
            rendered_pixel_width=width,
            rendered_pixel_height=height,
            crop_coordinates_px=(0, 0, width, height),
        )
        selections.append(
            NarrativeEvidenceSelection(
                reference,
                evidence,
                visual,
                "approved_narrative_support",
            )
        )
    return selections


def test_exactly_two_stable_template_ids_and_invalid_values_are_rejected() -> None:
    assert SUPPORTED_TEMPLATE_IDS == ("orange_bank_compact", "detailed_reference")
    payload = _request(_ids(1)).model_dump(mode="json")
    payload["template_id"] = "comparison_matrix"
    with pytest.raises(ValidationError):
        NarrativePresentationRequest.model_validate(payload)


def test_orange_source_range_numbering_and_roles_are_strict() -> None:
    mapper = OrangeBankCompactNarrativePptxGenerator.source_index
    assert ORANGE_SOURCE_FIRST_SLIDE == 10
    assert ORANGE_SOURCE_LAST_SLIDE == 29
    assert ORANGE_DIVIDER_SOURCE_SLIDE == 10
    assert ORANGE_COMPACT_SOURCE_SLIDES == tuple(range(11, 18))
    assert ORANGE_EVIDENCE_SOURCE_SLIDES == tuple(range(18, 30))
    assert mapper(10, "divider") == 0
    assert mapper(11, "compact_summary") == 1
    assert mapper(29, "evidence_attestation") == 19
    for outside in (9, 30):
        with pytest.raises(ValueError, match="outside PowerPoint 10-29"):
            mapper(outside, "compact_summary")
    with pytest.raises(ValueError, match="not eligible"):
        mapper(10, "compact_summary")
    with pytest.raises(ValueError, match="not eligible"):
        mapper(18, "divider")


@pytest.mark.parametrize(
    ("count", "expected_slides", "expected_summary_slides"),
    [(1, 3, 1), (3, 5, 1), (4, 7, 2)],
)
def test_compact_pagination_order_editability_and_evidence(
    tmp_path: Path,
    count: int,
    expected_slides: int,
    expected_summary_slides: int,
) -> None:
    ids = _ids(count)
    image_path = tmp_path / "approved-page.png"
    with fitz.open() as document:
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 96), "APPROVED SOURCE PAGE", fontsize=18)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(1, 1), alpha=False)
        pixmap.save(image_path)
    output = tmp_path / f"compact-{count}.pptx"
    request = _request(ids).model_copy(update={"template_id": "orange_bank_compact"})
    review = _review(ids)
    result = OrangeBankCompactNarrativePptxGenerator(Path(__file__).parents[1]).generate(
        output,
        request,
        review,
        _selections(ids, image_path),
    )

    deck = Presentation(output)
    with fitz.open(
        Path(__file__).parents[1]
        / "templates/reference_pack/source/OT_DVT_SDSI__OrangeBANK.pdf"
    ) as source_pdf:
        source_page = source_pdf[0].rect
    assert abs(deck.slide_width - round(source_page.width / 72 * 914400)) <= 1
    assert abs(deck.slide_height - round(source_page.height / 72 * 914400)) <= 1
    assert len(deck.slides) == expected_slides
    assert result.narrative_slide_count == 1 + expected_summary_slides
    assert result.evidence_slide_count == count
    assert [item["reference_id"] for item in result.reference_to_slide] == ids
    assert [item["reference_id"] for item in result.reference_to_evidence_slide] == ids
    if count == 4:
        assert [item["slide_number"] for item in result.reference_to_slide] == [2, 2, 2, 3]
        assert [item["card_index"] for item in result.reference_to_slide] == [1, 2, 3, 1]
    assert all(item["slide_number"] <= 1 + expected_summary_slides for item in result.reference_to_slide)
    assert all(item["slide_number"] > 1 + expected_summary_slides for item in result.reference_to_evidence_slide)
    assert sum(
        1
        for slide in deck.slides
        for shape in slide.shapes
        if shape.name == "Evidence source image"
    ) == count

    summary_shapes = {
        shape.name: shape
        for slide in list(deck.slides)[1 : 1 + expected_summary_slides]
        for shape in slide.shapes
    }
    for index in range(1, count + 1):
        assert summary_shapes[f"C.MISSION_{index}"].has_text_frame
        assert review.narrative.references[index - 1].headline.text in summary_shapes[f"C.MISSION_{index}"].text
        assert review.narrative.references[index - 1].realisations[0].text in summary_shapes[f"C.SERVICES_{index}"].text
        assert f"Trusted offering {index} / Trusted sector {index}" in summary_shapes[f"C.SERVICES_{index}"].text
        assert [line for line in summary_shapes[f"C.CLIENT_COUNTRY_{index}"].text.splitlines() if line] == [
            f"Trusted client {index}",
            f"Trusted country {index}",
        ]
    assert "C.MISSION_2" not in summary_shapes if count == 1 else True

    all_text = "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ).casefold()
    for forbidden in ("mewa", "mhrsd", "water sustainability", "point de contact", "placeholder"):
        assert forbidden not in all_text
    assert "…" not in all_text


def test_compact_keeps_reference_when_optional_evidence_is_unavailable(tmp_path: Path) -> None:
    ids = _ids(2)
    image_path = tmp_path / "approved-page.png"
    with fitz.open() as document:
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 96), "APPROVED SOURCE PAGE", fontsize=18)
        page.get_pixmap(matrix=fitz.Matrix(1, 1), alpha=False).save(image_path)
    all_selections = _selections(ids, image_path)
    output = tmp_path / "compact-optional-evidence.pptx"

    result = OrangeBankCompactNarrativePptxGenerator(Path(__file__).parents[1]).generate(
        output,
        _request(ids).model_copy(update={"template_id": "orange_bank_compact"}),
        _review(ids),
        all_selections[:1],
        trusted_references=[selection.reference for selection in all_selections],
    )

    assert [item["reference_id"] for item in result.reference_to_slide] == ids
    assert result.evidence_slide_count == 1
    assert result.reference_to_evidence_slide == [{"reference_id": ids[0], "slide_number": 3}]
    assert len(Presentation(output).slides) == 3


def test_compact_pdf_import_clones_source_geometry_removes_logos_and_keeps_editable_verdanatext(
    tmp_path: Path,
) -> None:
    root = Path(__file__).parents[1]
    generator = OrangeBankCompactNarrativePptxGenerator(root)
    source_hash_before = sha256_file(generator.source_path)
    ids = _ids(1)
    image_path = tmp_path / "approved-page.png"
    with fitz.open() as document:
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 96), "APPROVED SOURCE PAGE", fontsize=18)
        page.get_pixmap(matrix=fitz.Matrix(1, 1), alpha=False).save(image_path)
    output = tmp_path / "orange-cloned-source.pptx"
    generator.generate(
        output,
        _request(ids).model_copy(update={"template_id": "orange_bank_compact"}),
        _review(ids),
        _selections(ids, image_path),
    )
    assert sha256_file(generator.source_path) == source_hash_before == generator.source_sha256

    source = Presentation(generator.clone_base_path)
    generated = Presentation(output)
    assert len(source.slides) == 20
    assert (generated.slide_width, generated.slide_height) == (source.slide_width, source.slide_height)

    def geometry(shape) -> tuple[int, int, int, int, int]:
        return (int(shape.shape_type), shape.left, shape.top, shape.width, shape.height)

    divider_source = {
        geometry(shape)
        for shape in source.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
    }
    divider_generated = {
        geometry(shape)
        for shape in generated.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
    }
    assert divider_source <= divider_generated
    divider_artwork = max(
        (
            shape
            for shape in source.slides[0].shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ),
        key=lambda shape: shape.width * shape.height,
    )
    divider_artwork_hash = hashlib.sha256(divider_artwork.image.blob).hexdigest()
    assert divider_artwork_hash in {
        hashlib.sha256(shape.image.blob).hexdigest()
        for shape in generated.slides[0].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    }

    summary_source = {
        geometry(shape)
        for shape in source.slides[1].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
        or (shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.height / 914400 > 1.5)
    }
    summary_generated = {
        geometry(shape)
        for shape in generated.slides[1].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
        or (shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.height / 914400 > 1.5)
    }
    assert summary_source <= summary_generated

    evidence_source_frames = {
        geometry(shape)
        for shape in source.slides[8].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
    }
    evidence_generated_frames = {
        geometry(shape)
        for shape in generated.slides[2].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
    }
    assert evidence_source_frames <= evidence_generated_frames

    source_client_logo_hashes = {
        hashlib.sha256(shape.image.blob).hexdigest()
        for shape in source.slides[1].shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and shape.left / 914400 > 11
        and 1 < shape.top / 914400 < 6.7
        and 0.3 < shape.height / 914400 < 1.2
    }
    generated_picture_hashes = {
        hashlib.sha256(shape.image.blob).hexdigest()
        for slide in generated.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    }
    assert source_client_logo_hashes
    assert not (source_client_logo_hashes & generated_picture_hashes)

    editable_names = {
        shape.name
        for slide in generated.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    }
    assert {
        "C.MISSION_1",
        "C.SERVICES_1",
        "C.CLIENT_COUNTRY_1",
        "Evidence reference title",
    } <= editable_names
    font_names = {
        run.font.name
        for slide in generated.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip() and run.font.name
    }
    assert font_names == {"Verdana"}
    assert "Arial" not in font_names


def test_compact_overflow_blocks_without_truncating(tmp_path: Path) -> None:
    ids = _ids(1)
    image_path = Path(__file__).parents[1] / "templates/reference_pack/v1/assets/devoteam_logo.png"
    review = _review(ids)
    reference = review.narrative.references[0].model_copy(
        update={"realisations": [_text("Long approved compact delivery " * 220)]}
    )
    review = review.model_copy(
        update={"narrative": review.narrative.model_copy(update={"references": [reference]})}
    )
    output = tmp_path / "compact-overflow.pptx"
    with pytest.raises(PptxContentOverflowError) as captured:
        OrangeBankCompactNarrativePptxGenerator(Path(__file__).parents[1]).generate(
            output,
            _request(ids).model_copy(update={"template_id": "orange_bank_compact"}),
            review,
            _selections(ids, image_path),
        )
    assert captured.value.condition.reference_id == ids[0]
    assert captured.value.condition.field == "compact_services"
    assert not output.exists()


def test_compact_omits_overlong_optional_description_with_explicit_warning(tmp_path: Path) -> None:
    ids = _ids(1)
    image_path = Path(__file__).parents[1] / "templates/reference_pack/v1/assets/devoteam_logo.png"
    review = _review(ids)
    overlong = "Administrative approved source paragraph " * 8
    administrative_realisations = "Nous attestons par la présente que les prestations ont été réalisées. " * 4
    reference = review.narrative.references[0].model_copy(
        update={
            "short_description": _text(overlong),
            "realisations": [_text(administrative_realisations)],
        }
    )
    review = review.model_copy(
        update={"narrative": review.narrative.model_copy(update={"references": [reference]})}
    )
    output = tmp_path / "compact-optional-description.pptx"
    result = OrangeBankCompactNarrativePptxGenerator(Path(__file__).parents[1]).generate(
        output,
        _request(ids).model_copy(update={"template_id": "orange_bank_compact"}),
        review,
        _selections(ids, image_path),
    )
    assert result.export_warnings == [
        "COMPACT_SHORT_DESCRIPTION_OMITTED",
        "COMPACT_ADMINISTRATIVE_REALISATION_OMITTED",
    ]
    rendered_text = "\n".join(
        shape.text
        for slide in Presentation(output).slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert overlong not in rendered_text
    assert administrative_realisations not in rendered_text


def test_legacy_review_is_adapted_to_both_template_specific_schemas(tmp_path: Path) -> None:
    service, provider = _service(tmp_path)
    request = _service_request()
    detailed = service.generate(request.model_copy(update={"template_id": "detailed_reference"}))
    compact = service.generate(request.model_copy(update={"template_id": "orange_bank_compact"}))

    detailed_manifest = json.loads(
        (service.output_root / detailed.generation_id / "generation_manifest.json").read_text(encoding="utf-8")
    )
    compact_manifest = json.loads(
        (service.output_root / compact.generation_id / "generation_manifest.json").read_text(encoding="utf-8")
    )
    assert provider.calls == 0
    assert detailed_manifest["approved_narrative_status"] == compact_manifest["approved_narrative_status"] == "READY_FOR_PRESENTATION"
    assert detailed_manifest["selected_reference_ids"] == compact_manifest["selected_reference_ids"] == [REF_A]
    assert detailed_manifest["template_id"] == "detailed_reference"
    assert compact_manifest["template_id"] == "orange_bank_compact"
    assert compact_manifest["reference_slide_mappings"][0]["reference_id"] == REF_A
    assert compact_manifest["narrative_slide_mappings"][0]["slide_number"] == 1
    for response in (detailed, compact):
        directory = service.output_root / response.generation_id
        with fitz.open(directory / "narrative_reference_pack.pdf") as pdf:
            assert pdf.page_count == len(Presentation(directory / "narrative_reference_pack.pptx").slides)
