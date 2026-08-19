from __future__ import annotations

import json
from copy import deepcopy
from pathlib import Path

import fitz
import pandas as pd
import pytest
import yaml
from fastapi.testclient import TestClient
from pydantic import ValidationError
from pptx import Presentation

from app.api.dependencies import get_reference_pack_service
from app.api.main import app
from app.api.settings import load_config
from reference_pack.content_builder import prepare_reference
from reference_pack.pdf_converter import LibreOfficePdfConverter
from reference_pack.pptx_generator import PowerPointGenerator
from reference_pack.schemas import GenerationArtifacts, ReferencePackRequest, ReferencePackResponse, TrustedEvidence, TrustedReference
from reference_pack.service import ReferencePackService
from reference_pack.validation import ReferenceValidationError, TrustedV2Repository, sha256_file


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = yaml.safe_load((ROOT / "templates/reference_pack/v1/template_config.yaml").read_text(encoding="utf-8"))


def _id(character: str) -> str:
    return character * 64


def _trusted(index: int, *, title: str | None = None, language: str = "fr") -> TrustedReference:
    character = "abcdef0123456789"[index % 16]
    reference_id = _id(character)
    evidence = TrustedEvidence(
        chunk_id=_id("f" if character != "f" else "e"),
        document_id=f"doc-{index}",
        source_file_name=f"source-{index}.pdf",
        source_sha256=_id("1"),
        source_page=index + 1,
        citation_label=f"source-{index}.pdf — page {index + 1}",
        citation_uri=f"https://example.test/doc-{index}#page={index + 1}",
        language=language,
        display_text=(
            "Evidence confirms architecture, implementation and service-delivery activities for the selected project."
            if language != "ar"
            else "تؤكد الأدلة تنفيذ خدمات هندسة النظام وتطبيق المشروع وتسليم الخدمات للعميل."
        ),
    )
    mission = title or f"Mission de transformation numérique {index + 1}"
    return TrustedReference(
        reference_id=reference_id,
        row_number=index + 1,
        mission_title=mission,
        client=f"Client {index + 1}",
        country="Tunisie",
        period="2021–2023",
        sector="Banque",
        offering="Cloud",
        business_unit="Digital Impulse",
        description="Architecture du système. Mise en œuvre de la solution. Transfert de compétences.",
        services_delivered=["Architecture du système", "Mise en œuvre de la solution"],
        technologies=["Cloud"],
        capabilities=["Accompagnement à la mise en place"],
        evidence=[evidence],
    )


def _request(ids: list[str], *, language: str = "fr", formats: list[str] | None = None) -> ReferencePackRequest:
    return ReferencePackRequest(
        title="Références pertinentes pour la mission" if language == "fr" else "المراجع ذات الصلة بالفرصة",
        client_name="Client démonstration" if language == "fr" else "عميل تجريبي",
        subtitle="Sélection Devoteam" if language == "fr" else "مجموعة مراجع ديفوتيم",
        language=language,
        reference_ids=ids,
        output_formats=formats or ["pptx"],
    )


def _repository(retrieval: bool = True, display: bool = True, page: int = 1, eligible: bool = True) -> TrustedV2Repository:
    repo = TrustedV2Repository.__new__(TrustedV2Repository)
    repo.config = {"filters": {"allowed_security_classifications": ["INTERNAL"]}}
    repo.references = pd.DataFrame(
        [
            {
                "reference_id": _id("a"), "row_number": 1, "document_retrieval_eligible": eligible,
                "reference_number": "REF-1", "service_nature": "Architecture cloud et mise en œuvre",
                "offering": "Cloud", "client": "Client source", "country": "Tunisie",
                "project_year": "2022", "sector": "Banque", "business_unit": "Digital Impulse",
            }
        ]
    )
    repo.chunks = pd.DataFrame(
        [
            {
                "chunk_id": _id("b"), "document_id": "doc-1", "source_file_name": "source.pdf",
                "source_sha256": _id("c"), "page_number_1_based": page, "chunk_index_in_page": 0,
                "citation_label": "source.pdf — page 1", "citation_uri": "https://example.test/source#page=1",
                "reference_rows_json": "[1]", "approved_for_retrieval": retrieval,
                "approved_for_display": display, "display_text": "Project delivery evidence confirms cloud architecture and implementation services.",
                "security_classification": "INTERNAL", "page_language": "en", "document_language": "en",
                "source_relative_path": "raw/source.pdf",
            }
        ]
    )
    repo.quarantined = pd.DataFrame([{"reference_rows_json": "[1]"}]) if not eligible else pd.DataFrame(columns=["reference_rows_json"])
    return repo


def test_zero_and_duplicate_references_are_rejected() -> None:
    with pytest.raises(ValidationError):
        _request([])
    with pytest.raises(ValidationError, match="unique"):
        _request([_id("a"), _id("a")])


def test_presentation_metadata_is_sanitized_without_damaging_unicode() -> None:
    request = ReferencePackRequest(
        title="<script>alert('x')</script><b>Références عربية</b>",
        client_name="Client\x00 sûr",
        reference_ids=[_id("a")],
    )
    assert request.title == "Références عربية"
    assert request.client_name == "Client sûr"


def test_unknown_and_quarantined_reference_ids_are_rejected() -> None:
    repo = _repository()
    with pytest.raises(ReferenceValidationError) as unknown:
        repo.load_selected([_id("d")])
    assert unknown.value.reason == "UNKNOWN_REFERENCE_ID"
    with pytest.raises(ReferenceValidationError) as quarantined:
        _repository(eligible=False).load_selected([_id("a")])
    assert quarantined.value.reason == "REFERENCE_QUARANTINED"


def test_retrieval_only_display_prohibited_and_missing_pages_are_rejected() -> None:
    with pytest.raises(ReferenceValidationError) as retrieval_only:
        _repository(display=False).load_selected([_id("a")])
    assert retrieval_only.value.reason == "DISPLAY_EVIDENCE_REQUIRED"
    with pytest.raises(ReferenceValidationError) as prohibited:
        _repository(display=False).load_selected([_id("a")])
    assert prohibited.value.reason == "DISPLAY_EVIDENCE_REQUIRED"
    with pytest.raises(ReferenceValidationError) as missing_page:
        _repository(page=0).load_selected([_id("a")])
    assert missing_page.value.reason == "INVALID_SOURCE_LINEAGE"


def test_mission_title_stays_concise_instead_of_concatenating_service_scope() -> None:
    repo = _repository()
    repo.references.loc[0, "service_nature"] = (
        "Assistance à la transformation de la direction de l'organisation et des systèmes "
        "avec gouvernance, architecture, processus, sécurité, conduite du changement et formation"
    )
    reference = repo.load_selected([_id("a")])[0]
    assert len(reference.mission_title) <= 52
    assert "sécurité" not in reference.mission_title


@pytest.mark.parametrize("count,expected_slides", [(1, 5), (3, 9), (4, 12)])
def test_pptx_generation_and_summary_evidence_pagination(tmp_path: Path, count: int, expected_slides: int) -> None:
    trusted = [_trusted(index) for index in range(count)]
    prepared = [prepare_reference(reference, "fr") for reference in trusted]
    output = tmp_path / f"pack-{count}.pptx"
    result = PowerPointGenerator(ROOT, TEMPLATE).generate(output, _request([item.reference_id for item in trusted]), prepared)
    assert result.slide_count == expected_slides
    reopened = Presentation(output)
    assert len(reopened.slides) == expected_slides
    summary = [item for item in result.slide_provenance if item.slide_type == "reference_summary"]
    evidence = [item for item in result.slide_provenance if item.slide_type == "evidence_annex"]
    assert len(summary) == (count + 2) // 3
    assert len(evidence) == count
    assert all(len(item.reference_ids) <= 3 for item in summary)
    assert all(item.evidence_visuals[0]["fallback_reason"] for item in evidence)


def test_actual_approved_source_page_is_rendered_with_aspect_and_citation(tmp_path: Path) -> None:
    source_root = tmp_path / "sources"
    source = source_root / "raw/source.pdf"
    source.parent.mkdir(parents=True)
    with fitz.open() as document:
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 96), "APPROVED SOURCE PAGE", fontsize=18)
        page.insert_text((72, 140), "Cloud architecture and implementation evidence.", fontsize=12)
        document.save(source)

    reference = _trusted(0).model_copy(deep=True)
    reference.evidence[0] = reference.evidence[0].model_copy(
        update={
            "source_file_name": "source.pdf",
            "source_sha256": sha256_file(source),
            "source_page": 1,
            "source_relative_path": "raw/source.pdf",
            "citation_label": "source.pdf — page 1",
        }
    )
    template = deepcopy(TEMPLATE)
    template["footer"]["logo_path"] = str(
        (ROOT / TEMPLATE["footer"]["logo_path"]).resolve()
    )
    template["evidence"]["local_source_roots"] = ["sources"]
    output = tmp_path / "source-page.pptx"
    result = PowerPointGenerator(tmp_path, template).generate(
        output,
        _request([reference.reference_id]),
        [prepare_reference(reference, "fr")],
    )
    visual = result.evidence_visuals[0]
    assert visual["rendered_source_image"] is True
    assert visual["fallback_reason"] is None
    assert visual["aspect_ratio_preserved"] is True
    assert len(visual["crop_coordinates_px"]) == 4
    evidence_slide = next(
        slide for slide, provenance in zip(Presentation(output).slides, result.slide_provenance)
        if provenance.slide_type == "evidence_annex"
    )
    assert any(shape.name == "Evidence source image" for shape in evidence_slide.shapes)
    text = "\n".join(shape.text for shape in evidence_slide.shapes if hasattr(shape, "text"))
    assert "source.pdf" in text and "page 1" in text
    assert "Text evidence fallback" not in text


def test_stable_reference_order_and_missing_logo_fallback(tmp_path: Path) -> None:
    trusted = [_trusted(index) for index in (2, 0, 1)]
    prepared = [prepare_reference(reference, "fr") for reference in trusted]
    output = tmp_path / "ordered.pptx"
    result = PowerPointGenerator(ROOT, TEMPLATE).generate(output, _request([item.reference_id for item in trusted]), prepared)
    summary = next(item for item in result.slide_provenance if item.slide_type == "reference_summary")
    assert summary.reference_ids == [item.reference_id for item in trusted]
    text = "\n".join(shape.text for slide in Presentation(output).slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "Client 3" in text and "Client 1" in text and "Client 2" in text
    assert not (ROOT / "templates/reference_pack/v1/logo_registry.yaml").read_text(encoding="utf-8").count(".png")


def test_french_accents_arabic_unicode_and_no_internal_values_or_paths(tmp_path: Path) -> None:
    french = _trusted(0, title="Architecture de l’écosystème numérique")
    arabic = _trusted(1, title="هندسة النظام وتنفيذ الخدمات", language="ar")
    for language, reference in (("fr", french), ("ar", arabic)):
        output = tmp_path / f"unicode-{language}.pptx"
        PowerPointGenerator(ROOT, TEMPLATE).generate(output, _request([reference.reference_id], language=language), [prepare_reference(reference, language)])
        text = "\n".join(shape.text for slide in Presentation(output).slides for shape in slide.shapes if hasattr(shape, "text"))
        assert ("écosystème" in text) if language == "fr" else ("هندسة النظام" in text)
        lowered = text.casefold()
        assert "bm25" not in lowered and "dense cosine" not in lowered and "rrf" not in lowered
        assert "c:\\" not in lowered and "/users/" not in lowered


def test_pdf_conversion_page_count_and_unicode(tmp_path: Path) -> None:
    converter = LibreOfficePdfConverter(TEMPLATE["generation"]["libreoffice_candidates"])
    if converter.executable() is None:
        pytest.skip("LibreOffice is not installed")
    reference = _trusted(0, title="Architecture de l’écosystème numérique")
    pptx = tmp_path / "pdf-source.pptx"
    request = _request([reference.reference_id])
    generated = PowerPointGenerator(ROOT, TEMPLATE).generate(pptx, request, [prepare_reference(reference, "fr")])
    result = converter.convert(pptx, tmp_path / "converted.pdf", generated.slide_count, [request.title])
    assert result.warning is None and result.path is not None
    with fitz.open(result.path) as document:
        assert document.page_count == generated.slide_count


def test_real_manifest_completeness_hashes_and_path_traversal() -> None:
    service = ReferencePackService(ROOT, load_config())
    candidate_ids = service.repository.references.loc[
        service.repository.references["document_retrieval_eligible"], "reference_id"
    ].astype(str)
    selected = None
    for reference_id in candidate_ids:
        try:
            service.repository.load_selected([reference_id])
            selected = reference_id
            break
        except ReferenceValidationError:
            continue
    assert selected is not None
    artifact = service.generate(_request([selected]))
    manifest = artifact.manifest
    required = {
        "generation_id", "created_at_utc", "selected_reference_ids", "selected_ordering",
        "corpus_version", "corpus_manifest_sha256", "template_id", "template_version",
        "template_pdf_sha256", "source_documents", "source_pages", "evidence_chunk_ids",
        "evidence_visuals",
        "application_version", "outputs", "generation_warnings", "exact_generation_command",
    }
    assert required <= set(manifest)
    output = Path(artifact.directory) / "reference_pack.pptx"
    assert manifest["outputs"]["pptx"]["sha256"] == sha256_file(output)
    assert manifest["selected_reference_ids"] == [selected]
    assert manifest["corpus_version"] == "v2"
    with pytest.raises(ValueError):
        service.download_path("../escape", "pptx")


def test_reference_pack_api_contract_and_downloads(tmp_path: Path) -> None:
    generation_id = "reference-pack-20260803T120000000000Z-aaaaaaaaaa"
    for name, content in (("reference_pack.pptx", b"pptx"), ("reference_pack.pdf", b"%PDF-demo"), ("generation_manifest.json", b"{}")):
        (tmp_path / name).write_bytes(content)

    response = ReferencePackResponse(
        generation_id=generation_id,
        status="completed",
        selected_reference_count=1,
        slide_count=5,
        pptx_download_url=f"/api/reference-packs/{generation_id}/download/pptx",
        pdf_download_url=f"/api/reference-packs/{generation_id}/download/pdf",
        manifest_download_url=f"/api/reference-packs/{generation_id}/download/manifest",
    )

    class FakeReferencePackService:
        def generate(self, request: ReferencePackRequest) -> GenerationArtifacts:
            assert request.reference_ids == [_id("a")]
            return GenerationArtifacts(response=response, directory=str(tmp_path), manifest={})

        def get(self, requested_id: str) -> ReferencePackResponse:
            assert requested_id == generation_id
            return response

        def download_path(self, requested_id: str, kind: str) -> Path:
            assert requested_id == generation_id
            return tmp_path / {"pptx": "reference_pack.pptx", "pdf": "reference_pack.pdf", "manifest": "generation_manifest.json"}[kind]

    app.dependency_overrides[get_reference_pack_service] = lambda: FakeReferencePackService()
    client = TestClient(app)
    try:
        created = client.post("/api/reference-packs", json=_request([_id("a")]).model_dump(mode="json"))
        assert created.status_code == 201
        assert created.json()["generation_id"] == generation_id
        status = client.get(f"/api/reference-packs/{generation_id}")
        assert status.status_code == 200 and status.json()["slide_count"] == 5
        for kind in ("pptx", "pdf", "manifest"):
            download = client.get(f"/api/reference-packs/{generation_id}/download/{kind}")
            assert download.status_code == 200
    finally:
        app.dependency_overrides.pop(get_reference_pack_service, None)
