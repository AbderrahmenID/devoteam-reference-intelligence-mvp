from __future__ import annotations

import hashlib
import json
import zipfile
from pathlib import Path
from types import SimpleNamespace

import pytest
import fitz
import yaml
from fastapi.testclient import TestClient
from pptx import Presentation

from app.api.dependencies import get_narrative_presentation_service
from app.api.main import app
from app.api.settings import PROJECT_ROOT, load_config
from reference_narrative.pptx_generator import NarrativePptxGenerator
from reference_narrative.presentation_schemas import NarrativePresentationRequest
from reference_narrative.presentation_service import (
    NarrativePresentationExportError,
    NarrativePresentationService,
)
from reference_narrative.schemas import (
    EditableReferenceSectionNarrative,
    FieldSupportPlan,
    NarrativeGenerationRequest,
    NarrativeReferenceMetadata,
    NarrativeReviewResponse,
    NarrativeSupportPlan,
    NarrativeValidationResult,
    ReferenceNarrative,
    ReferenceNarrativeDraft,
    ReferenceSectionNarrative,
    SectionSupportPlan,
    SupportedDetailedPresentationCopy,
    SupportedDetailedRealisation,
    SupportedNarrativeText,
)
from reference_narrative.service import ReferenceNarrativeService
from reference_narrative.template_mapper import PptxContentOverflowError
from reference_pack.evidence_renderer import EvidenceRenderer
from reference_pack.validation import sha256_file

from test_reference_narrative import FakeRepository, REF_A, _reference_a


TEMPLATE = PROJECT_ROOT / "templates/reference_pack/source/references sapmple and template.pptx"
SAMPLE_MEDIA_SHA1 = {
    "51939cfa864b25a3fd815b8d0494953f47e58c14",  # sample landscape
    "964cfe8064569051db76509de55ea7d63dc4729c",  # sample flag
    "90edb490e65f5ad7e96f8c854122b31c15df3c00",  # sample client logo
}


def _text(value: str) -> SupportedNarrativeText:
    return SupportedNarrativeText(text=value, support_ids=[])


def _ids(count: int) -> list[str]:
    return [character * 64 for character in "abcd"[:count]]


def _editable(ids: list[str], *, challenge: str = "A grounded challenge.") -> EditableReferenceSectionNarrative:
    return EditableReferenceSectionNarrative(
        section_intro="Grounded section introduction.",
        overall_storyline="Grounded overall storyline.",
        why_these_references="Grounded reason for this selection.",
        references=[
            ReferenceNarrativeDraft(
                headline=f"Approved headline {index}",
                challenge=challenge,
                realisations=[f"Approved delivery {index}.1", f"Approved delivery {index}.2"],
                benefits=[f"Approved benefit {index}.1"],
            )
            for index, _ in enumerate(ids, start=1)
        ],
    )


def _request(
    ids: list[str],
    *,
    language: str = "fr",
    approved: bool = True,
    status: str = "READY_FOR_PRESENTATION",
    narrative: EditableReferenceSectionNarrative | None = None,
) -> NarrativePresentationRequest:
    return NarrativePresentationRequest(
        generation_request=NarrativeGenerationRequest(
            selected_reference_ids=ids,
            opportunity_title="Reference opportunity",
            target_language=language,
        ),
        narrative=narrative or _editable(ids),
        approved=approved,
        approved_narrative_status=status,
        approved_reference_ids=ids,
    )


def _review(
    ids: list[str],
    *,
    language: str = "fr",
    empty_challenge: bool = False,
    empty_benefits: bool = False,
    long_challenge: bool = False,
) -> NarrativeReviewResponse:
    localized = {
        "fr": ("Référence approuvée", "Réalisation vérifiée", "Bénéfice vérifié"),
        "en": ("Approved reference", "Verified delivery", "Verified benefit"),
        "ar": ("مرجع معتمد", "إنجاز موثق", "فائدة موثقة"),
    }[language]
    references = []
    metadata = []
    for index, reference_id in enumerate(ids, start=1):
        challenge = "" if empty_challenge else "A grounded operational challenge."
        if long_challenge:
            challenge = "Overflowing reviewed challenge content " * 200
        benefits = [] if empty_benefits else [_text(f"{localized[2]} {index}")]
        references.append(
            ReferenceNarrative(
                reference_id=reference_id,
                headline=_text(f"{localized[0]} {index}"),
                short_description=_text(""),
                challenge=_text(challenge),
                devoteam_contribution=_text(""),
                realisations=[_text(f"{localized[1]} {index}.1"), _text(f"{localized[1]} {index}.2")],
                benefits=benefits,
                why_relevant_to_opportunity=_text(""),
                detailed_presentation=SupportedDetailedPresentationCopy(
                    mission_title=_text(f"{localized[0]} {index}"),
                    challenges=[] if empty_challenge else [_text(challenge)],
                    realisations=[
                        SupportedDetailedRealisation(
                            text=_text(f"{localized[1]} {index}.1"),
                            subitems=[],
                        ),
                        SupportedDetailedRealisation(
                            text=_text(f"{localized[1]} {index}.2"),
                            subitems=[],
                        ),
                    ],
                    benefits=benefits,
                ),
            )
        )
        metadata.append(
            NarrativeReferenceMetadata(
                reference_id=reference_id,
                mission_title=f"Trusted mission {index}",
                client=f"Trusted client {index}",
                country=f"Trusted country {index}",
                sector=f"Trusted sector {index}",
                period=f"20{20 + index}",
                offering=f"Trusted offering {index}",
            )
        )
    section_text = {
        "fr": ("Introduction approuvée.", "Fil conducteur approuvé.", "Sélection approuvée."),
        "en": ("Approved introduction.", "Approved overall storyline.", "Approved selection rationale."),
        "ar": ("مقدمة معتمدة.", "سرد عام معتمد.", "مبررات اختيار معتمدة."),
    }[language]
    return NarrativeReviewResponse(
        narrative=ReferenceSectionNarrative(
            section_intro=_text(section_text[0]),
            overall_storyline=_text(section_text[1]),
            why_these_references=_text(section_text[2]),
            references=references,
        ),
        validation=NarrativeValidationResult(valid=True, export_blocked=False, warnings=[]),
        warnings=[],
        support_plan=NarrativeSupportPlan(
            references=[FieldSupportPlan(reference_id=reference_id) for reference_id in ids],
            section=SectionSupportPlan(),
        ),
        reference_metadata=metadata,
    )


def _all_text(path: Path) -> str:
    return "\n".join(
        shape.text
        for slide in Presentation(path).slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )


def _shape(slide, name: str):
    def walk(shapes):
        for shape in shapes:
            if shape.name == name:
                return shape
            if hasattr(shape, "shapes"):
                found = walk(shape.shapes)
                if found is not None:
                    return found
        return None

    found = walk(slide.shapes)
    if found is None:
        raise AssertionError(f"Shape {name} was not found")
    return found


def test_one_and_four_references_create_one_slide_per_reference_and_preserve_order(tmp_path: Path) -> None:
    generator = NarrativePptxGenerator(PROJECT_ROOT)
    for count, expected in ((1, 1), (4, 4)):
        ids = _ids(count)
        output = tmp_path / f"{count}.pptx"
        result = generator.generate(output, _request(ids), _review(ids))
        reopened = Presentation(output)
        assert len(reopened.slides) == expected
        assert result.reference_to_slide == [
            {"reference_id": reference_id, "slide_number": index + 1}
            for index, reference_id in enumerate(ids)
        ]
        for index in range(count):
            assert f"Trusted client {index + 1}" in _all_text(output)
            assert _shape(reopened.slides[index], "D.CLIENT").text == f"Trusted client {index + 1}"


def test_trusted_metadata_and_approved_case_content_are_editable_native_text(tmp_path: Path) -> None:
    output = tmp_path / "editable.pptx"
    NarrativePptxGenerator(PROJECT_ROOT).generate(output, _request(_ids(1)), _review(_ids(1)))
    slide = Presentation(output).slides[0]
    assert _shape(slide, "D.CLIENT").text == "Trusted client 1"
    assert _shape(slide, "D.SECTOR").text == "Trusted sector 1"
    assert _shape(slide, "D.PERIOD").text == "2021"
    assert "Référence approuvée 1" in _shape(slide, "D.MISSION_TITLE").text
    assert "A grounded operational challenge." in _shape(slide, "D.CHALLENGE").text
    assert "Réalisation vérifiée 1.1" in _shape(slide, "D.REALISATIONS").text
    assert "Bénéfice vérifié 1" in _shape(slide, "D.BENEFITS").text
    bullet_paragraphs = _shape(slide, "D.REALISATIONS").text_frame.paragraphs[1:]
    assert [paragraph.text for paragraph in bullet_paragraphs] == [
        "Réalisation vérifiée 1.1",
        "Réalisation vérifiée 1.2",
    ]
    assert all(paragraph._p.xpath("./a:pPr/a:buChar") for paragraph in bullet_paragraphs)


def test_empty_fields_keep_required_headings_without_inventing_content(tmp_path: Path) -> None:
    output = tmp_path / "empty.pptx"
    review = _review(_ids(1), empty_challenge=True, empty_benefits=True)
    NarrativePptxGenerator(PROJECT_ROOT).generate(output, _request(_ids(1)), review)
    names = {shape.name for shape in Presentation(output).slides[0].shapes}
    assert "D.CHALLENGE" in names
    assert "D.BENEFITS" in names
    assert _shape(Presentation(output).slides[0], "D.CHALLENGE").text.strip() == "Challenges"
    assert _shape(Presentation(output).slides[0], "D.BENEFITS").text.strip() == "Bénéfices"
    assert "challenge unavailable" not in _all_text(output).casefold()
    assert "benefit unavailable" not in _all_text(output).casefold()


def test_sample_text_logos_and_template_annotations_do_not_leak(tmp_path: Path) -> None:
    output = tmp_path / "safe.pptx"
    NarrativePptxGenerator(PROJECT_ROOT).generate(output, _request(_ids(1)), _review(_ids(1)))
    text = _all_text(output).casefold()
    for forbidden in (
        "mewa",
        "mhrsd",
        "water sustainability",
        "secteur public",
        "accelerator",
        "strategic domain",
        "point de contact",
        "dispositif",
    ):
        assert forbidden not in text
    with zipfile.ZipFile(output) as archive:
        media_hashes = {
            hashlib.sha1(archive.read(name)).hexdigest()
            for name in archive.namelist()
            if name.startswith("ppt/media/")
        }
        package_text = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist()
            if name.endswith((".xml", ".rels", ".txt"))
        ).casefold()
    assert not (media_hashes & SAMPLE_MEDIA_SHA1)
    for forbidden in ("mewa", "mhrsd", "water sustainability", "secteur public", "point de contact"):
        assert forbidden not in package_text


def test_detailed_output_preserves_source_slide_dimensions(tmp_path: Path) -> None:
    output = tmp_path / "source-geometry.pptx"
    NarrativePptxGenerator(PROJECT_ROOT).generate(output, _request(_ids(1)), _review(_ids(1)))
    source = Presentation(
        PROJECT_ROOT / "templates/reference_pack/source/references sapmple and template.pptx"
    )
    generated = Presentation(output)
    assert generated.slide_width == source.slide_width
    assert generated.slide_height == source.slide_height


@pytest.mark.parametrize(
    ("language", "expected"),
    [("fr", "Référence approuvée"), ("en", "Approved reference"), ("ar", "مرجع معتمد")],
)
def test_french_english_and_arabic_unicode_survive_pptx_structure(
    tmp_path: Path, language: str, expected: str
) -> None:
    output = tmp_path / f"{language}.pptx"
    ids = _ids(1)
    NarrativePptxGenerator(PROJECT_ROOT).generate(
        output,
        _request(ids, language=language),
        _review(ids, language=language),
    )
    assert expected in _all_text(output)
    assert len(Presentation(output).slides) == 1


def test_overflow_blocks_generation_without_truncation(tmp_path: Path) -> None:
    output = tmp_path / "overflow.pptx"
    ids = _ids(1)
    with pytest.raises(PptxContentOverflowError) as captured:
        NarrativePptxGenerator(PROJECT_ROOT).generate(
            output,
            _request(ids),
            _review(ids, long_challenge=True),
        )
    assert captured.value.condition.reference_id == ids[0]
    assert captured.value.condition.field == "challenge"
    assert not output.exists()


class ExplodingProvider:
    provider_name = "must-not-run"
    model_name = "must-not-run"

    def __init__(self):
        self.calls = 0

    def generate(self, messages, response_schema):
        self.calls += 1
        raise AssertionError("The LLM provider must not be called during PPTX export")


def _service(tmp_path: Path) -> tuple[NarrativePresentationService, ExplodingProvider]:
    reference, document_types = _reference_a()
    fixture_root = PROJECT_ROOT / ".tmp/pytest-narrative-pptx" / tmp_path.name / "trusted_sources"
    fixture_root.mkdir(parents=True, exist_ok=True)
    source_path = fixture_root / "source.pdf"
    document = fitz.open()
    page = document.new_page(width=595, height=842)
    page.insert_text((72, 96), "Approved BCT evidence page", fontsize=18)
    document.save(source_path)
    document.close()
    reference.evidence[0] = reference.evidence[0].model_copy(
        update={
            "source_sha256": sha256_file(source_path),
            "source_relative_path": "source.pdf",
        }
    )
    repository = FakeRepository([reference], document_types)
    repository.identity = SimpleNamespace(version="test-v2")
    provider = ExplodingProvider()
    validator = ReferenceNarrativeService(repository, provider)
    render_config = yaml.safe_load(
        (PROJECT_ROOT / "templates/reference_pack/v1/template_config.yaml").read_text(encoding="utf-8")
    )
    render_config["evidence"]["local_source_roots"] = [
        fixture_root.relative_to(PROJECT_ROOT).as_posix()
    ]
    service = NarrativePresentationService(
        PROJECT_ROOT,
        load_config(),
        repository=repository,
        validator=validator,
        evidence_renderer=EvidenceRenderer(PROJECT_ROOT, render_config),
    )
    service.output_root = PROJECT_ROOT / ".tmp/pytest-narrative-pptx" / tmp_path.name
    service.output_root.mkdir(parents=True, exist_ok=True)
    return service, provider


def _service_request(*, approved: bool = True, status: str = "READY_FOR_PRESENTATION", contribution: str = "") -> NarrativePresentationRequest:
    narrative = EditableReferenceSectionNarrative(
        references=[
            ReferenceNarrativeDraft(
                headline="Operationalisation du PCA de la BCT",
                devoteam_contribution=contribution,
            )
        ]
    )
    return NarrativePresentationRequest(
        generation_request=NarrativeGenerationRequest(
            selected_reference_ids=[REF_A],
            opportunity_title="PCA opportunity",
            target_language="fr",
        ),
        narrative=narrative,
        approved=approved,
        approved_narrative_status=status,
        approved_reference_ids=[REF_A],
    )


def test_export_gate_rejects_unapproved_and_blocking_narratives(tmp_path: Path) -> None:
    service, provider = _service(tmp_path)
    with pytest.raises(NarrativePresentationExportError) as unapproved:
        service.generate(_service_request(approved=False, status="DRAFT"))
    assert unapproved.value.reason == "NARRATIVE_NOT_APPROVED"

    changed = _service_request().model_copy(update={"approved_reference_ids": ["b" * 64]})
    with pytest.raises(NarrativePresentationExportError) as changed_set:
        service.generate(changed)
    assert changed_set.value.reason == "NARRATIVE_REFERENCE_SET_CHANGED"

    with pytest.raises(NarrativePresentationExportError) as blocking:
        service.generate(_service_request(contribution="The project delivered a positive ROI."))
    assert blocking.value.reason == "NARRATIVE_HAS_BLOCKING_WARNINGS"
    assert provider.calls == 0


def test_service_creates_manifest_hashes_and_never_calls_model(tmp_path: Path) -> None:
    template_hash_before = sha256_file(TEMPLATE)
    service, provider = _service(tmp_path)
    response = service.generate(_service_request())
    directory = service.output_root / response.generation_id
    manifest = json.loads((directory / "generation_manifest.json").read_text(encoding="utf-8"))
    reviewed = directory / "reviewed_content.json"
    assert response.slide_count == 1
    assert (directory / "narrative_reference_pack.pptx").is_file()
    assert (directory / "narrative_reference_pack.pdf").is_file()
    assert manifest["template_source_sha256"] == template_hash_before
    assert manifest["reviewed_content_sha256"] == hashlib.sha256(
        json.dumps(
            json.loads(reviewed.read_text(encoding="utf-8")),
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    assert manifest["reference_to_slide"] == [{"reference_id": REF_A, "slide_number": 1}]
    assert manifest["reference_to_evidence_slide"] == []
    assert manifest["narrative_slide_count"] == 1
    assert manifest["evidence_slide_count"] == 0
    assert manifest["outputs"]["pdf"]["sha256"] == sha256_file(directory / "narrative_reference_pack.pdf")
    assert manifest["font_substitution"]["used"] == "Montserrat / Montserrat Light"
    assert manifest["overflow_validation"]["status"] == "PASS"
    assert provider.calls == 0
    assert sha256_file(TEMPLATE) == template_hash_before
    assert len(Presentation(directory / "narrative_reference_pack.pptx").slides) == 1


def test_api_generation_and_download_contract(tmp_path: Path) -> None:
    service, _provider = _service(tmp_path)
    app.dependency_overrides[get_narrative_presentation_service] = lambda: service
    try:
        client = TestClient(app)
        created = client.post(
            "/api/reference-narrative/presentations",
            json=_service_request().model_dump(mode="json"),
        )
        assert created.status_code == 201
        body = created.json()
        downloaded = client.get(body["pptx_download_url"])
        pdf = client.get(body["pdf_download_url"])
        manifest = client.get(body["manifest_download_url"])
    finally:
        app.dependency_overrides.pop(get_narrative_presentation_service, None)
    assert downloaded.status_code == 200
    assert downloaded.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert pdf.status_code == 200
    assert pdf.headers["content-type"].startswith("application/pdf")
    assert manifest.status_code == 200
