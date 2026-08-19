from __future__ import annotations

import io
import json
from pathlib import Path

import fitz
import pytest
from PIL import Image
from fastapi.testclient import TestClient
from pptx import Presentation

from app.api.dependencies import get_narrative_presentation_service
from app.api.main import app
from app.api.settings import PROJECT_ROOT
from reference_narrative.evidence_annex import NarrativeEvidenceSelection, choose_evidence
from reference_narrative.compact_pptx_generator import OrangeBankCompactNarrativePptxGenerator
from reference_narrative.pptx_generator import NarrativePptxGenerator
from reference_narrative.presentation_service import NarrativePresentationExportError
from reference_pack.evidence_renderer import EvidenceVisual
from reference_pack.pdf_converter import PdfConversionResult
from reference_pack.schemas import TrustedEvidence, TrustedReference
from reference_pack.validation import ReferenceValidationError, sha256_file

from test_narrative_pptx import REF_A, _ids, _request, _review, _service, _service_request, _shape


def _visible_text(path: Path) -> str:
    return "\n".join(
        shape.text
        for slide in Presentation(path).slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def test_complete_one_reference_pack_contains_real_page_editable_narrative_and_matching_pdf(tmp_path: Path) -> None:
    service, provider = _service(tmp_path)
    response = service.generate(
        _service_request().model_copy(update={"template_id": "orange_bank_compact"})
    )
    directory = service.output_root / response.generation_id
    pptx_path = directory / "narrative_reference_pack.pptx"
    pdf_path = directory / "narrative_reference_pack.pdf"
    manifest = json.loads((directory / "generation_manifest.json").read_text(encoding="utf-8"))

    deck = Presentation(pptx_path)
    assert len(deck.slides) == 3
    assert _shape(deck.slides[1], "C.MISSION_1").has_text_frame
    picture = _shape(deck.slides[2], "Evidence source image")
    with Image.open(io.BytesIO(picture.image.blob)) as image:
        embedded_ratio = image.width / image.height
    assert abs(embedded_ratio - picture.width / picture.height) <= 0.01
    evidence_text = "\n".join(
        shape.text for shape in deck.slides[2].shapes if getattr(shape, "has_text_frame", False)
    )
    assert "BCT attestation.pdf" in evidence_text
    assert "page 1" in evidence_text
    assert "Approved BCT evidence page" not in evidence_text
    for forbidden in ("chunk-a", "retrieval score", "raw\\evidence", str(PROJECT_ROOT)):
        assert forbidden.casefold() not in evidence_text.casefold()

    with fitz.open(pdf_path) as pdf:
        assert pdf.page_count == len(deck.slides)
    assert manifest["narrative_slide_count"] == 2
    assert manifest["evidence_slide_count"] == 1
    assert manifest["total_slide_count"] == 3
    assert manifest["reference_to_evidence_slide"] == [{"reference_id": REF_A, "slide_number": 3}]
    assert manifest["evidence_pages"][0]["rendered_source_image"] is True
    assert manifest["evidence_pages"][0]["aspect_ratio_preserved"] is True
    assert manifest["evidence_pages"][0]["source_hash_validation"] == "PASS"
    assert manifest["outputs"]["pptx"]["sha256"] == sha256_file(pptx_path)
    assert manifest["outputs"]["pdf"]["sha256"] == sha256_file(pdf_path)
    assert "source_relative_path" not in json.dumps(manifest)
    assert provider.calls == 0


def test_four_reference_annex_order_is_after_all_narrative_slides(tmp_path: Path) -> None:
    source = PROJECT_ROOT / "templates/reference_pack/source/OT_DVT_SDSI__OrangeBANK.pdf"
    image_path = tmp_path / "page.png"
    with fitz.open(source) as document:
        pixmap = document[0].get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        pixmap.save(image_path)
    with Image.open(image_path) as image:
        width, height = image.size

    ids = _ids(4)
    selections: list[NarrativeEvidenceSelection] = []
    for index, reference_id in enumerate(ids, start=1):
        evidence = TrustedEvidence(
            chunk_id=f"chunk-{index}",
            document_id=f"document-{index}",
            source_file_name=f"approved-source-{index}.pdf",
            source_sha256=sha256_file(source),
            source_page=1,
            citation_label=f"Approved source {index} - page 1",
            citation_uri="https://example.test/source",
            language="fr",
            display_text=f"Approved evidence {index}",
        )
        reference = TrustedReference(
            reference_id=reference_id,
            row_number=index,
            mission_title=f"Trusted mission {index}",
            client=f"Trusted client {index}",
            country="Tunisia",
            period="2024",
            sector="Technology",
            offering="Cloud",
            business_unit="Cloud",
            description="Approved reference",
            services_delivered=[],
            technologies=[],
            capabilities=[],
            evidence=[evidence],
        )
        visual = EvidenceVisual(
            reference_id=reference_id,
            chunk_id=evidence.chunk_id,
            document_id=evidence.document_id,
            source_file_name=evidence.source_file_name,
            source_page=1,
            source_sha256=evidence.source_sha256,
            image_path=image_path,
            rendering_method="pdf_page_render",
            fallback_reason=None,
            original_pixel_width=width,
            original_pixel_height=height,
            rendered_pixel_width=width,
            rendered_pixel_height=height,
            crop_coordinates_px=(0, 0, width, height),
        )
        selections.append(NarrativeEvidenceSelection(reference, evidence, visual, "approved_narrative_support"))

    output = tmp_path / "four-reference-pack.pptx"
    result = OrangeBankCompactNarrativePptxGenerator(PROJECT_ROOT).generate(
        output,
        _request(ids).model_copy(update={"template_id": "orange_bank_compact"}),
        _review(ids),
        selections,
    )
    assert result.slide_count == 7
    assert result.reference_to_evidence_slide == [
        {"reference_id": reference_id, "slide_number": index + 4}
        for index, reference_id in enumerate(ids)
    ]
    deck = Presentation(output)
    assert [_shape(deck.slides[index], "Evidence source label").text for index in range(3, 7)] == [
        f"Source: approved-source-{index}.pdf · page 1" for index in range(1, 5)
    ]


@pytest.mark.parametrize(
    ("language", "client", "expected_label"),
    [
        ("fr", "Banque Éthique", "Justificatifs de nos références"),
        ("en", "Ethical Bank", "Reference evidence"),
        ("ar", "البنك المركزي", "أدلة المراجع"),
    ],
)
def test_evidence_labels_preserve_french_english_and_arabic_unicode(
    tmp_path: Path, language: str, client: str, expected_label: str
) -> None:
    image_path = PROJECT_ROOT / "templates/reference_pack/v1/assets/devoteam_logo.png"
    with Image.open(image_path) as image:
        width, height = image.size
    reference_id = "a" * 64
    evidence = TrustedEvidence(
        chunk_id="unicode-chunk",
        document_id="unicode-document",
        source_file_name="preuve-approuvée.pdf",
        source_sha256="1" * 64,
        source_page=1,
        citation_label="Approved page 1",
        citation_uri="https://example.test/source",
        language=language,
        display_text="Approved evidence",
    )
    reference = TrustedReference(
        reference_id=reference_id,
        row_number=1,
        mission_title="مهمة موثقة" if language == "ar" else "Mission approuvée",
        client=client,
        country="Tunisia",
        period="2024",
        sector="Banking",
        offering="Cloud",
        business_unit="Cloud",
        description="Approved reference",
        services_delivered=[],
        technologies=[],
        capabilities=[],
        evidence=[evidence],
    )
    visual = EvidenceVisual(
        reference_id=reference_id,
        chunk_id=evidence.chunk_id,
        document_id=evidence.document_id,
        source_file_name=evidence.source_file_name,
        source_page=1,
        source_sha256=evidence.source_sha256,
        image_path=image_path,
        rendering_method="source_image_render",
        fallback_reason=None,
        original_pixel_width=width,
        original_pixel_height=height,
        rendered_pixel_width=width,
        rendered_pixel_height=height,
        crop_coordinates_px=(0, 0, width, height),
    )
    output = tmp_path / f"unicode-annex-{language}.pptx"
    OrangeBankCompactNarrativePptxGenerator(PROJECT_ROOT).generate(
        output,
        _request([reference_id], language=language).model_copy(update={"template_id": "orange_bank_compact"}),
        _review([reference_id], language=language),
        [NarrativeEvidenceSelection(reference, evidence, visual, "highest_priority_display_evidence")],
    )
    text = _visible_text(output)
    assert expected_label in text
    assert client in text
    assert "preuve-approuvée.pdf" in text


@pytest.mark.parametrize(
    ("update", "reason"),
    [
        ({"source_relative_path": "missing.pdf"}, "EVIDENCE_SOURCE_NOT_FOUND"),
        ({"source_sha256": "0" * 64}, "EVIDENCE_HASH_MISMATCH"),
        ({"source_page": 99}, "EVIDENCE_PAGE_NOT_APPROVED"),
    ],
)
def test_required_evidence_failures_are_explicit_and_leave_no_generation(
    tmp_path: Path, update: dict[str, object], reason: str
) -> None:
    service, _provider = _service(tmp_path)
    reference = service.repository.by_id[REF_A]
    reference.evidence[0] = reference.evidence[0].model_copy(update=update)
    before = set(service.output_root.glob("narrative-pptx-*"))
    with pytest.raises(NarrativePresentationExportError) as captured:
        service.generate(
            _service_request().model_copy(update={"template_id": "orange_bank_compact"})
        )
    assert captured.value.reason == reason
    assert set(service.output_root.glob("narrative-pptx-*")) == before


def test_quarantined_and_retrieval_only_repositories_are_rejected(tmp_path: Path) -> None:
    for reason in ("REFERENCE_QUARANTINED", "DISPLAY_EVIDENCE_REQUIRED"):
        service, _provider = _service(tmp_path)

        def reject(_ids, rejected_reason=reason):
            raise ReferenceValidationError(rejected_reason, [REF_A], "Evidence is not display eligible")

        service.repository.load_selected = reject
        with pytest.raises(ReferenceValidationError) as captured:
            service.generate(_service_request())
        assert captured.value.reason == reason


def test_browser_cannot_inject_evidence_path(tmp_path: Path) -> None:
    service, _provider = _service(tmp_path)
    app.dependency_overrides[get_narrative_presentation_service] = lambda: service
    payload = _service_request().model_dump(mode="json")
    payload["source_path"] = r"C:\untrusted\source.pdf"
    try:
        response = TestClient(app).post("/api/reference-narrative/presentations", json=payload)
    finally:
        app.dependency_overrides.pop(get_narrative_presentation_service, None)
    assert response.status_code == 422
    assert "source_path" in response.text


def test_wrong_reference_evidence_cannot_be_selected() -> None:
    service_review = _review(["b" * 64])
    service, _provider = _service(Path("wrong-reference-evidence"))
    reference = service.repository.by_id[REF_A]
    with pytest.raises(ValueError, match="order"):
        choose_evidence([reference], service_review, {})


def test_pdf_conversion_failure_is_blocking(tmp_path: Path) -> None:
    service, _provider = _service(tmp_path)

    class FailingConverter:
        def convert(self, *_args, **_kwargs):
            return PdfConversionResult(
                path=None,
                warning="conversion failed",
                validation={"status": "FAIL"},
                command=[],
            )

    service.pdf_converter = FailingConverter()
    with pytest.raises(NarrativePresentationExportError) as captured:
        service.generate(_service_request())
    assert captured.value.reason == "PDF_CONVERSION_FAILED"
    assert not list(service.output_root.glob("narrative-pptx-*"))
