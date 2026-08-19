from __future__ import annotations

import copy
import math
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import yaml
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400
CORAL = RGBColor(0xFF, 0x49, 0x6B)
BODY = RGBColor(0x16, 0x16, 0x16)
MUTED = RGBColor(0x55, 0x55, 0x59)


@dataclass(frozen=True)
class OverflowCondition:
    reference_id: str | None
    field: str
    required_lines: int
    available_lines: int
    minimum_font_pt: int


@dataclass(frozen=True)
class TextFitMeasurement:
    fits: bool
    font_size_pt: int
    required_lines: int
    available_lines: int


class PptxContentOverflowError(RuntimeError):
    def __init__(self, condition: OverflowCondition):
        super().__init__(
            f"PPTX_CONTENT_OVERFLOW: {condition.field} requires approximately "
            f"{condition.required_lines} lines but the template allows {condition.available_lines} "
            f"at {condition.minimum_font_pt} pt"
        )
        self.condition = condition


def load_template_mapping(project_root: Path) -> tuple[dict[str, Any], dict[str, Any]]:
    registry_path = project_root / "templates/reference_pack/qwen_studio/template_registry.yaml"
    registry = yaml.safe_load(registry_path.read_text(encoding="utf-8"))
    entry = registry["templates"]["detailed_reference"]
    mapping = yaml.safe_load((project_root / entry["mapping_file"]).read_text(encoding="utf-8"))
    return entry, mapping


def duplicate_source_slide(presentation: Presentation, source_slide_index: int):
    """Duplicate an inherited source slide without flattening editable objects."""
    source = presentation.slides[source_slide_index]
    destination = presentation.slides.add_slide(source.slide_layout)
    destination.notes_slide

    destination_layout = next(
        relationship
        for relationship in destination.part.rels.values()
        if relationship.reltype.endswith("/slideLayout")
    )
    destination_notes = next(
        relationship
        for relationship in destination.part.rels.values()
        if relationship.reltype.endswith("/notesSlide")
    )
    relationship_ids: dict[str, str] = {}
    for relationship in source.part.rels.values():
        if relationship.reltype.endswith("/slideLayout"):
            relationship_ids[relationship.rId] = destination_layout.rId
            continue
        if relationship.reltype.endswith("/notesSlide"):
            relationship_ids[relationship.rId] = destination_notes.rId
            continue
        relationship_ids[relationship.rId] = destination.part.rels._add_relationship(
            relationship.reltype,
            relationship.target_ref if relationship.is_external else relationship.target_part,
            relationship.is_external,
        )

    destination_element = destination._element
    for child in list(destination_element):
        destination_element.remove(child)
    for child in source._element:
        destination_element.append(copy.deepcopy(child))
    destination_element.attrib.clear()
    destination_element.attrib.update(source._element.attrib)
    for element in destination_element.iter():
        for attribute, value in list(element.attrib.items()):
            if value in relationship_ids:
                element.set(attribute, relationship_ids[value])
    destination.__dict__.pop("shapes", None)
    return destination


def remove_original_slides(presentation: Presentation, count: int) -> None:
    for _ in range(count):
        slide_id = presentation.slides._sldIdLst[0]
        presentation.part.drop_rel(slide_id.rId)
        del presentation.slides._sldIdLst[0]


def find_shape(slide, shape_id: int):
    def walk(shapes):
        for shape in shapes:
            if shape.shape_id == shape_id:
                return shape
            if hasattr(shape, "shapes"):
                found = walk(shape.shapes)
                if found is not None:
                    return found
        return None

    found = walk(slide.shapes)
    if found is None:
        raise RuntimeError(f"Template shape ID {shape_id} was not found")
    return found


def delete_shape(slide, shape_id: int) -> None:
    shape = find_shape(slide, shape_id)
    shape._element.getparent().remove(shape._element)


def replace_with_transparent_textbox(slide, shape, name: str):
    """Rebuild a template placeholder as an editable, transparent text box."""
    box = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
    box.name = name
    box.fill.background()
    box.line.fill.background()
    shape._element.getparent().remove(shape._element)
    return box


def cleanup_unused_image_relationships(slide) -> None:
    used = set(slide._element.xpath(".//a:blip/@r:embed"))
    for relationship in list(slide.part.rels.values()):
        if relationship.reltype.endswith("/image") and relationship.rId not in used:
            slide.part.drop_rel(relationship.rId)


def _remove_bullet_nodes(paragraph) -> None:
    properties = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone", "a:buFont", "a:buSzPct"):
        for node in properties.findall(qn(tag)):
            properties.remove(node)


def _clear_bullet(paragraph) -> None:
    _remove_bullet_nodes(paragraph)
    paragraph._p.get_or_add_pPr().append(etree.Element(qn("a:buNone")))


def _set_native_bullet(paragraph, *, level: int = 0, font_name: str = "Montserrat") -> None:
    _remove_bullet_nodes(paragraph)
    properties = paragraph._p.get_or_add_pPr()
    properties.set("lvl", str(level))
    if level == 0:
        properties.set("marL", "180000")
        properties.set("indent", "-152400")
        bullet_character = "●"
    else:
        properties.set("marL", "914400")
        properties.set("indent", "-285750")
        bullet_character = "○"
    bullet_font = etree.Element(qn("a:buFont"))
    bullet_font.set("typeface", font_name)
    properties.append(bullet_font)
    marker = etree.Element(qn("a:buChar"))
    marker.set("char", bullet_character)
    properties.append(marker)


def _format_paragraph(paragraph, font_name: str, font_size: int, *, bold: bool, color: RGBColor, rtl: bool) -> None:
    paragraph.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(1)
    paragraph.line_spacing = 1.0
    if rtl:
        paragraph._p.get_or_add_pPr().set("rtl", "1")
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def _wrapped_line_count(text: str, characters_per_line: int) -> int:
    lines = 0
    for raw_line in str(text).splitlines() or [""]:
        wrapped = textwrap.wrap(raw_line.strip(), width=characters_per_line, break_long_words=True)
        lines += max(1, len(wrapped))
    return lines


def measure_text_fit(
    *,
    heading: str,
    values: list[str],
    width_inches: float,
    height_inches: float,
    intended_pt: int,
    minimum_pt: int,
    calibrated_lines: int | None = None,
) -> TextFitMeasurement:
    """Measure copy with the exact deterministic rules used by the renderers."""
    last_required = 0
    last_capacity = 0
    for size in range(intended_pt, minimum_pt - 1, -1):
        characters_per_line = max(12, math.floor(width_inches * 72 / (size * 0.55)))
        required = _wrapped_line_count(heading, characters_per_line) if heading else 0
        required += sum(_wrapped_line_count(value, characters_per_line - 2) for value in values)
        capacity = max(
            1,
            math.floor(height_inches * 72 / (size * 1.12)),
            int(calibrated_lines or 0),
        )
        last_required, last_capacity = required, capacity
        if required <= capacity:
            return TextFitMeasurement(True, size, required, capacity)
    return TextFitMeasurement(False, minimum_pt, last_required, last_capacity)


def choose_font_size(
    *,
    heading: str,
    values: list[str],
    width_inches: float,
    height_inches: float,
    intended_pt: int,
    minimum_pt: int,
    reference_id: str | None,
    field: str,
    calibrated_lines: int | None = None,
) -> int:
    measurement = measure_text_fit(
        heading=heading,
        values=values,
        width_inches=width_inches,
        height_inches=height_inches,
        intended_pt=intended_pt,
        minimum_pt=minimum_pt,
        calibrated_lines=calibrated_lines,
    )
    if measurement.fits:
        return measurement.font_size_pt
    raise PptxContentOverflowError(
        OverflowCondition(
            reference_id,
            field,
            measurement.required_lines,
            measurement.available_lines,
            minimum_pt,
        )
    )


def set_text_zone(
    shape,
    *,
    heading: str,
    values: list[str],
    font_name: str,
    body_font_name: str | None = None,
    intended_pt: int,
    minimum_pt: int,
    reference_id: str | None,
    field: str,
    bullets: bool,
    rtl: bool,
    calibrated_lines: int | None = None,
) -> int:
    clean_values = [value.strip() for value in values if value and value.strip()]
    width = shape.width / EMU_PER_INCH
    height = shape.height / EMU_PER_INCH
    size = choose_font_size(
        heading=heading,
        values=clean_values,
        width_inches=width,
        height_inches=height,
        intended_pt=intended_pt,
        minimum_pt=minimum_pt,
        reference_id=reference_id,
        field=field,
        calibrated_lines=calibrated_lines,
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = None
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    header = frame.paragraphs[0]
    header.text = heading
    _clear_bullet(header)
    _format_paragraph(header, font_name, intended_pt, bold=True, color=CORAL, rtl=rtl)
    body_font = body_font_name or font_name
    for value in clean_values:
        paragraph = frame.add_paragraph()
        paragraph.text = value
        if bullets:
            _set_native_bullet(paragraph)
        else:
            _clear_bullet(paragraph)
        _format_paragraph(paragraph, body_font, size, bold=False, color=BODY, rtl=rtl)
    return size


def set_hierarchical_text_zone(
    shape,
    *,
    heading: str,
    values,
    font_name: str,
    body_font_name: str | None = None,
    intended_pt: int,
    minimum_pt: int,
    reference_id: str | None,
    field: str,
    rtl: bool,
    calibrated_lines: int | None = None,
) -> int:
    """Render editable PowerPoint parent/child bullets without simulated spaces."""
    clean_values = []
    fit_values: list[str] = []
    for item in values:
        text = str(item.get("text", "") if isinstance(item, dict) else item.text).strip()
        raw_subitems = item.get("subitems", []) if isinstance(item, dict) else item.subitems
        subitems = [str(value).strip() for value in raw_subitems if str(value).strip()]
        if text:
            clean_values.append((text, subitems))
            fit_values.extend([text, *[f"○ {subitem}" for subitem in subitems]])
    size = choose_font_size(
        heading=heading,
        values=fit_values,
        width_inches=shape.width / EMU_PER_INCH,
        height_inches=shape.height / EMU_PER_INCH,
        intended_pt=intended_pt,
        minimum_pt=minimum_pt,
        reference_id=reference_id,
        field=field,
        calibrated_lines=calibrated_lines,
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = None
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    header = frame.paragraphs[0]
    header.text = heading
    _clear_bullet(header)
    _format_paragraph(header, font_name, intended_pt, bold=True, color=CORAL, rtl=rtl)
    body_font = body_font_name or font_name
    for text, subitems in clean_values:
        paragraph = frame.add_paragraph()
        paragraph.text = text
        paragraph.level = 0
        _set_native_bullet(paragraph, level=0, font_name=body_font)
        _format_paragraph(paragraph, body_font, size, bold=False, color=BODY, rtl=rtl)
        for subitem in subitems:
            paragraph = frame.add_paragraph()
            paragraph.text = subitem
            paragraph.level = 1
            _set_native_bullet(paragraph, level=1, font_name=body_font)
            _format_paragraph(paragraph, body_font, size, bold=False, color=BODY, rtl=rtl)
    return size


def set_plain_text(shape, text: str, *, font_name: str, size: int, bold: bool = False, color: RGBColor = BODY, rtl: bool = False) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = None
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text.strip()
    _clear_bullet(paragraph)
    _format_paragraph(paragraph, font_name, size, bold=bold, color=color, rtl=rtl)


def set_fitted_plain_text(
    shape,
    text: str,
    *,
    font_name: str,
    intended_pt: int,
    minimum_pt: int,
    reference_id: str | None,
    field: str,
    bold: bool = False,
    color: RGBColor = BODY,
    rtl: bool = False,
) -> int:
    size = choose_font_size(
        heading="",
        values=[text],
        width_inches=shape.width / EMU_PER_INCH,
        height_inches=shape.height / EMU_PER_INCH,
        intended_pt=intended_pt,
        minimum_pt=minimum_pt,
        reference_id=reference_id,
        field=field,
    )
    set_plain_text(shape, text, font_name=font_name, size=size, bold=bold, color=color, rtl=rtl)
    return size


def add_plain_text(slide, name: str, zone: list[float], text: str, *, font_name: str, size: int, bold: bool = False, color: RGBColor = BODY, rtl: bool = False):
    box = slide.shapes.add_textbox(*(Inches(value) for value in zone))
    box.name = name
    set_plain_text(box, text, font_name=font_name, size=size, bold=bold, color=color, rtl=rtl)
    return box


def set_notes(slide, lines: list[str]) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.text = "\n".join(["[Sources]", *lines])


def clear_notes(slide) -> None:
    slide.notes_slide.notes_text_frame.text = ""
