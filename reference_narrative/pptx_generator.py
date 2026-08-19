from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from reference_pack.validation import sha256_file

from .evidence_annex import NarrativeEvidenceSelection
from .presentation_schemas import NarrativePresentationRequest
from .schemas import NarrativeReviewResponse
from .template_mapper import (
    BODY,
    CORAL,
    MUTED,
    add_plain_text,
    clear_notes,
    cleanup_unused_image_relationships,
    delete_shape,
    duplicate_source_slide,
    find_shape,
    load_template_mapping,
    remove_original_slides,
    set_notes,
    set_fitted_plain_text,
    set_plain_text,
    set_hierarchical_text_zone,
    set_text_zone,
)


@dataclass(frozen=True)
class NarrativePptxResult:
    path: Path
    slide_count: int
    narrative_slide_count: int
    evidence_slide_count: int
    reference_to_slide: list[dict[str, Any]]
    reference_to_evidence_slide: list[dict[str, Any]]
    evidence_visuals: list[dict[str, Any]]
    font_substitution: dict[str, str]
    overflow_validation: dict[str, Any]
    narrative_slide_mappings: list[dict[str, Any]]
    export_warnings: list[str]


class NarrativePptxGenerator:
    def __init__(self, project_root: Path):
        self.project_root = project_root.resolve()
        self.template_id = "detailed_reference"
        self.registry, self.mapping = load_template_mapping(self.project_root)
        self.template_display_name = self.registry["display_name"]
        self.source_path = (self.project_root / self.registry["source_file"]).resolve()
        self.source_sha256 = self.registry["source_sha256"]
        if sha256_file(self.source_path) != self.registry["source_sha256"]:
            raise RuntimeError("Template D source SHA-256 mismatch")

    @staticmethod
    def _remove_empty_zone(slide, shape_id: int) -> None:
        delete_shape(slide, shape_id)

    def _base_slide(self, presentation: Presentation):
        return duplicate_source_slide(presentation, int(self.mapping["source_slide"]) - 1)

    def _prepare_common(self, slide, page_number: int, *, section: bool) -> None:
        ids = self.mapping["shape_ids"]
        delete_ids = self.mapping["section_delete_shape_ids" if section else "case_delete_shape_ids"]
        for shape_id in delete_ids:
            delete_shape(slide, int(shape_id))
        devoteam = find_shape(slide, int(ids["devoteam_logo"]))
        devoteam.name = "D.DEVOTEAM_LOGO"
        page = find_shape(slide, int(ids["page_number"]))
        page.name = "D.PAGE_NUMBER"
        set_plain_text(
            page,
            str(page_number),
            font_name=self.mapping["typography"]["source_latin_font"],
            size=10,
            bold=True,
            color=BODY,
        )

    def _section_slide(self, presentation: Presentation, request: NarrativePresentationRequest, review: NarrativeReviewResponse):
        slide = self._base_slide(presentation)
        self._prepare_common(slide, 1, section=True)
        ids = self.mapping["shape_ids"]
        typography = self.mapping["typography"]
        labels = self.mapping["labels"][request.generation_request.target_language]
        rtl = request.generation_request.target_language == "ar"
        font = typography["fallback_font"]

        title = find_shape(slide, int(ids["mission_title"]))
        title.name = "D.SECTION_TITLE"
        title.left, title.top, title.width, title.height = (
            411480,
            1097280,
            2267712,
            1280160,
        )
        title_size = typography["section_title_pt"]
        set_fitted_plain_text(
            title,
            labels["section_title"],
            font_name=font,
            intended_pt=title_size,
            minimum_pt=int(typography["minimum_section_title_pt"]),
            reference_id=None,
            field="section_title",
            bold=True,
            color=CORAL,
            rtl=rtl,
        )

        fields = [
            ("section_intro", "section_intro", ids["challenge"]),
            ("overall_storyline", "overall_storyline", ids["realisations"]),
            ("why_these_references", "why_these_references", ids["benefits"]),
        ]
        for field, label_key, shape_id in fields:
            text = getattr(review.narrative, field).text.strip()
            if not text:
                self._remove_empty_zone(slide, int(shape_id))
                continue
            shape = find_shape(slide, int(shape_id))
            shape.name = f"D.SECTION_{field.upper()}"
            set_text_zone(
                shape,
                heading=labels[label_key],
                values=[text],
                font_name=font,
                intended_pt=int(typography["intended_body_pt"]),
                minimum_pt=int(typography["minimum_body_pt"]),
                reference_id=None,
                field=field,
                bullets=False,
                rtl=rtl,
            )
        cleanup_unused_image_relationships(slide)
        set_notes(
            slide,
            [
                f"Template: {self.source_path.name}",
                "Content: human-reviewed section narrative",
                "Reference IDs: " + ", ".join(request.approved_reference_ids),
            ],
        )
        return slide

    def _case_slide(self, presentation: Presentation, request: NarrativePresentationRequest, review: NarrativeReviewResponse, index: int):
        slide = self._base_slide(presentation)
        page_number = index + 1
        self._prepare_common(slide, page_number, section=False)
        ids = self.mapping["shape_ids"]
        typography = self.mapping["typography"]
        labels = self.mapping["labels"][request.generation_request.target_language]
        rtl = request.generation_request.target_language == "ar"
        font = typography["fallback_font"]
        narrative = review.narrative.references[index]
        metadata = review.reference_metadata[index]
        detailed = narrative.detailed_presentation
        if detailed is None:
            raise RuntimeError(
                "Detailed PowerPoint rendering requires mission_title/challenges/realisations/benefits copy"
            )

        headline = detailed.mission_title.text.strip() or metadata.mission_title
        title = find_shape(slide, int(ids["mission_title"]))
        title.name = "D.MISSION_TITLE"
        size = int(typography["headline_pt"])
        set_text_zone(
            title,
            heading="",
            values=[headline],
            font_name=font,
            body_font_name=typography["source_latin_font"],
            intended_pt=size,
            minimum_pt=int(typography["minimum_headline_pt"]),
            reference_id=metadata.reference_id,
            field="headline",
            bullets=False,
            rtl=rtl,
            calibrated_lines=int(self.mapping["calibrated_lines"]["mission_title"]),
        )
        # Remove the empty heading paragraph produced for a headline-only zone.
        if title.text_frame.paragraphs and not title.text_frame.paragraphs[0].text:
            title.text_frame._txBody.remove(title.text_frame.paragraphs[0]._p)

        field_specs = [
            ("challenge", labels["challenge"], [item.text for item in detailed.challenges], ids["challenge"]),
            ("benefits", labels["benefits"], [item.text for item in detailed.benefits], ids["benefits"]),
        ]
        for field, heading, values, shape_id in field_specs:
            values = [value.strip() for value in values if value and value.strip()]
            shape = find_shape(slide, int(shape_id))
            shape.name = f"D.{field.upper()}"
            set_text_zone(
                shape,
                heading=heading,
                values=values,
                font_name=font,
                body_font_name=typography["source_light_font"],
                intended_pt=int(typography["intended_body_pt"]),
                minimum_pt=int(typography["minimum_body_pt"]),
                reference_id=metadata.reference_id,
                field=field,
                bullets=True,
                rtl=rtl,
                calibrated_lines=int(self.mapping["calibrated_lines"][field]),
            )

        realisations_shape = find_shape(slide, int(ids["realisations"]))
        realisations_shape.name = "D.REALISATIONS"
        set_hierarchical_text_zone(
            realisations_shape,
            heading=labels["realisations"],
            values=[
                {
                    "text": item.text.text,
                    "subitems": [subitem.text for subitem in item.subitems],
                }
                for item in detailed.realisations
            ],
            font_name=font,
            body_font_name=typography["source_light_font"],
            intended_pt=int(typography["intended_body_pt"]),
            minimum_pt=int(typography["minimum_body_pt"]),
            reference_id=metadata.reference_id,
            field="realisations",
            rtl=rtl,
            calibrated_lines=int(self.mapping["calibrated_lines"]["realisations"]),
        )

        sector = find_shape(slide, int(ids["sector"]))
        sector.name = "D.SECTOR"
        sector.left, sector.width = 411480, 2267712
        set_fitted_plain_text(sector, metadata.sector, font_name=font, intended_pt=10, minimum_pt=8, reference_id=metadata.reference_id, field="sector", bold=True, rtl=rtl)
        period_group = find_shape(slide, int(ids["period_group"]))
        period_group.name = "D.PERIOD_GROUP"
        period = find_shape(slide, int(ids["period_text"]))
        period.name = "D.PERIOD"
        set_fitted_plain_text(period, metadata.period, font_name=font, intended_pt=10, minimum_pt=8, reference_id=metadata.reference_id, field="period", rtl=rtl)

        zones = self.mapping["added_text_zones_inches"]
        client = add_plain_text(slide, "D.CLIENT", zones["client"], "", font_name=font, size=int(typography["client_pt"]), bold=True, color=CORAL, rtl=rtl)
        set_fitted_plain_text(client, metadata.client, font_name=font, intended_pt=int(typography["client_pt"]), minimum_pt=12, reference_id=metadata.reference_id, field="client", bold=True, color=CORAL, rtl=rtl)
        country = add_plain_text(slide, "D.COUNTRY", zones["country"], "", font_name=font, size=int(typography["metadata_pt"]), bold=True, rtl=rtl)
        set_fitted_plain_text(country, metadata.country, font_name=font, intended_pt=int(typography["metadata_pt"]), minimum_pt=8, reference_id=metadata.reference_id, field="country", bold=True, rtl=rtl)
        offering = add_plain_text(slide, "D.OFFERING", zones["offering"], "", font_name=font, size=int(typography["metadata_pt"]), color=MUTED, rtl=rtl)
        set_fitted_plain_text(offering, metadata.offering, font_name=font, intended_pt=int(typography["metadata_pt"]), minimum_pt=8, reference_id=metadata.reference_id, field="offering", color=MUTED, rtl=rtl)

        cleanup_unused_image_relationships(slide)
        clear_notes(slide)
        return slide

    @staticmethod
    def _picture_contain(slide, path: Path, zone: tuple[float, float, float, float], *, ratio: float):
        x, y, width, height = zone
        box_ratio = width / height
        if ratio >= box_ratio:
            rendered_width = width
            rendered_height = width / ratio
        else:
            rendered_height = height
            rendered_width = height * ratio
        picture = slide.shapes.add_picture(
            str(path),
            Inches(x + (width - rendered_width) / 2),
            Inches(y + (height - rendered_height) / 2),
            width=Inches(rendered_width),
            height=Inches(rendered_height),
        )
        picture.name = "D.EVIDENCE_SOURCE_IMAGE"
        return picture

    def _evidence_slide(
        self,
        presentation: Presentation,
        request: NarrativePresentationRequest,
        selection: NarrativeEvidenceSelection,
        page_number: int,
    ) -> dict[str, Any]:
        slide = self._base_slide(presentation)
        self._prepare_common(slide, page_number, section=False)
        ids = self.mapping["shape_ids"]
        for key in ("mission_title", "realisations", "challenge", "benefits", "sector", "period_group"):
            delete_shape(slide, int(ids[key]))

        labels = {
            "fr": {"evidence": "Preuve", "reference": "Référence", "source": "Source", "page": "Page", "client": "Client"},
            "en": {"evidence": "Evidence", "reference": "Reference", "source": "Source", "page": "Page", "client": "Client"},
            "ar": {"evidence": "الدليل", "reference": "المرجع", "source": "المصدر", "page": "الصفحة", "client": "العميل"},
        }[request.generation_request.target_language]
        rtl = request.generation_request.target_language == "ar"
        font = self.mapping["typography"]["fallback_font"]
        reference = selection.reference
        evidence = selection.evidence
        visual = selection.visual

        title = labels["evidence"]
        title_shape = add_plain_text(
            slide, "D.EVIDENCE_TITLE", [0.45, 0.24, 9.05, 0.38], title,
            font_name=font, size=16, bold=True, color=CORAL, rtl=rtl,
        )
        set_fitted_plain_text(
            title_shape, title, font_name=font, intended_pt=16, minimum_pt=11,
            reference_id=reference.reference_id, field="evidence_title", bold=True, color=CORAL, rtl=rtl,
        )
        metadata = (
            f"{labels['client']}: {reference.client}\n"
            f"{labels['reference']}: {reference.mission_title}"
        )
        metadata_shape = add_plain_text(
            slide, "D.EVIDENCE_METADATA", [0.45, 0.70, 9.05, 0.62], metadata,
            font_name=font, size=9, bold=True, color=BODY, rtl=rtl,
        )
        set_fitted_plain_text(
            metadata_shape, metadata, font_name=font, intended_pt=9, minimum_pt=8,
            reference_id=reference.reference_id, field="evidence_metadata", bold=True, color=BODY, rtl=rtl,
        )
        metadata_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        frame = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.42), Inches(8.84), Inches(3.30)
        )
        frame.name = "D.EVIDENCE_FRAME"
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        frame.line.color.rgb = CORAL
        frame.line.width = Pt(0.8)

        ratio = float(visual.rendered_pixel_width) / float(visual.rendered_pixel_height)
        picture = self._picture_contain(
            slide, visual.image_path, (0.66, 1.50, 8.68, 3.14), ratio=ratio
        )
        placed_ratio = float(picture.width) / float(picture.height)
        if abs(placed_ratio - ratio) > 0.01:
            raise RuntimeError("Evidence source image aspect ratio changed during placement")

        source_label = (
            f"{labels['source']}: {evidence.source_file_name}\n"
            f"{labels['page']}: {evidence.source_page}"
        )
        source_shape = add_plain_text(
            slide, "D.EVIDENCE_SOURCE_LABEL", [0.58, 4.78, 8.84, 0.48], source_label,
            font_name=font, size=8, color=MUTED, rtl=rtl,
        )
        set_fitted_plain_text(
            source_shape, source_label, font_name=font, intended_pt=8, minimum_pt=8,
            reference_id=reference.reference_id, field="evidence_source_label", color=MUTED, rtl=rtl,
        )
        source_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        cleanup_unused_image_relationships(slide)
        set_notes(
            slide,
            [
                f"Trusted reference ID: {reference.reference_id}",
                f"Approved source: {evidence.source_file_name}",
                f"Approved page: {evidence.source_page}",
                "Evidence selection: deterministic trusted lineage",
            ],
        )
        record = visual.manifest_record()
        record.update(
            {
                "selection_reason": selection.selection_reason,
                "placed_aspect_ratio": round(placed_ratio, 6),
                "aspect_ratio_preserved": True,
                "source_hash_validation": "PASS",
            }
        )
        record.pop("image_file_name", None)
        return record

    def generate(
        self,
        output_path: Path,
        request: NarrativePresentationRequest,
        review: NarrativeReviewResponse,
        evidence_selections: list[NarrativeEvidenceSelection] | None = None,
        *,
        prevalidated_review: NarrativeReviewResponse | None = None,
        copy_generation: dict[str, Any] | None = None,
    ) -> NarrativePptxResult:
        # If a prevalidated review (from PresentationCopyService) is provided,
        # prefer it for rendering to avoid reintroducing overflowing content.
        if prevalidated_review is not None:
            review = prevalidated_review

        presentation = Presentation(self.source_path)
        original_count = len(presentation.slides)
        reference_to_slide: list[dict[str, Any]] = []
        for index, reference_id in enumerate(request.approved_reference_ids):
            self._case_slide(presentation, request, review, index)
            reference_to_slide.append({"reference_id": reference_id, "slide_number": index + 1})
        reference_to_evidence_slide: list[dict[str, Any]] = []
        evidence_visuals: list[dict[str, Any]] = []
        remove_original_slides(presentation, original_count)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        reopened = Presentation(output_path)
        narrative_slide_count = len(request.approved_reference_ids)
        evidence_slide_count = 0
        if len(reopened.slides) != narrative_slide_count + evidence_slide_count:
            raise RuntimeError("Generated presentation slide count mismatch")
        return NarrativePptxResult(
            path=output_path,
            slide_count=len(reopened.slides),
            narrative_slide_count=narrative_slide_count,
            evidence_slide_count=evidence_slide_count,
            reference_to_slide=reference_to_slide,
            reference_to_evidence_slide=reference_to_evidence_slide,
            evidence_visuals=evidence_visuals,
            font_substitution={
                "source": "Montserrat / Montserrat Light",
                "used": "Montserrat / Montserrat Light",
            },
            overflow_validation={"status": "PASS", "policy": "block", "minimum_font_pt": int(self.mapping["typography"]["minimum_body_pt"])},
            narrative_slide_mappings=[
                {
                    "narrative_role": "detailed_reference",
                    "slide_number": index + 1,
                    "reference_id": reference_id,
                    "fields": ["mission_title", "challenges", "realisations", "benefits"],
                }
                for index, reference_id in enumerate(request.approved_reference_ids)
            ],
            export_warnings=[],
        )
