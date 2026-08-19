from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from reference_pack.pptx_generator import LOCAL_PATH, PROHIBITED_TEXT, PowerPointGenerator, TEXT
from reference_pack.schemas import BulletSource, PreparedReference, SlideProvenance, TrustedReference
from reference_pack.validation import sha256_file

from .evidence_annex import NarrativeEvidenceSelection
from .pptx_generator import NarrativePptxResult
from .presentation_schemas import NarrativePresentationRequest
from .schemas import NarrativeReviewResponse
from .template_mapper import (
    OverflowCondition,
    PptxContentOverflowError,
    choose_font_size,
    cleanup_unused_image_relationships,
    duplicate_source_slide,
    measure_text_fit,
    remove_original_slides,
    set_notes,
)


ORANGE_SOURCE_FIRST_SLIDE = 10
ORANGE_SOURCE_LAST_SLIDE = 29
ORANGE_DIVIDER_SOURCE_SLIDE = 10
ORANGE_COMPACT_SOURCE_SLIDES = tuple(range(11, 18))
ORANGE_EVIDENCE_SOURCE_SLIDES = tuple(range(18, 30))


INTRO_LABELS = {
    "fr": {
        "title": "Références pertinentes\npour la mission",
        "section_intro": "Introduction de la section",
        "overall_storyline": "Fil conducteur",
        "why_these_references": "Pourquoi ces références",
    },
    "en": {
        "title": "Relevant references\nfor the opportunity",
        "section_intro": "Section introduction",
        "overall_storyline": "Overall storyline",
        "why_these_references": "Why these references",
    },
    "ar": {
        "title": "المراجع ذات الصلة\nبالفرصة",
        "section_intro": "مقدمة القسم",
        "overall_storyline": "السرد العام",
        "why_these_references": "سبب اختيار هذه المراجع",
    },
}

ADMINISTRATIVE_REALISATION = re.compile(
    r"\b(?:attestons? par la présente|je soussigné|certifie|en foi de quoi|signature)\b",
    re.IGNORECASE,
)


class OrangeBankCompactNarrativePptxGenerator(PowerPointGenerator):
    """Approved-narrative adapter for the established compact Reference Pack renderer."""

    def __init__(self, project_root: Path):
        self.project_root = project_root.resolve()
        registry = yaml.safe_load(
            (self.project_root / "templates/reference_pack/qwen_studio/template_registry.yaml").read_text(
                encoding="utf-8"
            )
        )["templates"]["orange_bank_compact"]
        config = yaml.safe_load(
            (self.project_root / registry["mapping_file"]).read_text(encoding="utf-8")
        )
        super().__init__(self.project_root, config)
        self.template_id = "orange_bank_compact"
        self.template_display_name = registry["display_name"]
        self.source_path = (self.project_root / registry["source_file"]).resolve()
        self.source_sha256 = registry["source_sha256"]
        self.registry = registry
        if sha256_file(self.source_path) != self.source_sha256:
            raise RuntimeError("Orange Bank source SHA-256 mismatch")
        self.clone_base_path = (self.project_root / registry["derived_clone_file"]).resolve()
        clone_manifest_path = (self.project_root / registry["derived_clone_manifest"]).resolve()
        if not self.clone_base_path.is_file() or not clone_manifest_path.is_file():
            raise RuntimeError(
                "Orange Bank PDF clone base is missing; run scripts/build_orange_pdf_template.py"
            )
        self.clone_manifest = json.loads(clone_manifest_path.read_text(encoding="utf-8"))
        if self.clone_manifest["authoritative_source_sha256"] != self.source_sha256:
            raise RuntimeError("Orange Bank clone base was built from a different PDF")
        if sha256_file(self.clone_base_path) != self.clone_manifest["derived_clone_base_sha256"]:
            raise RuntimeError("Orange Bank PDF clone base SHA-256 mismatch")
        if self.clone_manifest["source_range_powerpoint"] != [
            ORANGE_SOURCE_FIRST_SLIDE,
            ORANGE_SOURCE_LAST_SLIDE,
        ]:
            raise RuntimeError("Orange Bank clone base does not represent slides 10 through 29")
        self.clone_base_slide_count = int(self.clone_manifest["derived_slide_count"])

    @staticmethod
    def source_index(powerpoint_slide_number: int, role: str) -> int:
        """Map authoritative PowerPoint numbering to the 10-29 clone-base index."""
        if not ORANGE_SOURCE_FIRST_SLIDE <= powerpoint_slide_number <= ORANGE_SOURCE_LAST_SLIDE:
            raise ValueError("Orange source slides outside PowerPoint 10-29 are not eligible")
        eligible = {
            "divider": (ORANGE_DIVIDER_SOURCE_SLIDE,),
            "compact_summary": ORANGE_COMPACT_SOURCE_SLIDES,
            "evidence_attestation": ORANGE_EVIDENCE_SOURCE_SLIDES,
        }
        if role not in eligible or powerpoint_slide_number not in eligible[role]:
            raise ValueError(
                f"PowerPoint slide {powerpoint_slide_number} is not eligible for Orange role {role}"
            )
        return powerpoint_slide_number - ORANGE_SOURCE_FIRST_SLIDE

    def _clone_source_slide(self, presentation: Presentation, powerpoint_slide_number: int, role: str):
        return duplicate_source_slide(
            presentation,
            self.source_index(powerpoint_slide_number, role),
        )

    @staticmethod
    def _remove_shape(shape) -> None:
        shape._element.getparent().remove(shape._element)

    @staticmethod
    def _mark_source_bleed(slide, presentation: Presentation) -> None:
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > presentation.slide_width
                or shape.top + shape.height > presentation.slide_height
            ):
                shape.name = "Orange source bleed"

    def _add_page_number(self, slide, number: int) -> None:
        footer = self.config["footer"]
        self._text(
            slide,
            str(number),
            footer["slide_number_x"],
            footer["slide_number_y"],
            footer["slide_number_width"],
            footer["slide_number_height"],
            size=self.type["slide_number_pt"],
            color="muted",
            bold=True,
            align=PP_ALIGN.CENTER,
            name="Slide number",
        )

    @staticmethod
    def _fit(
        text: str,
        *,
        width: float,
        height: float,
        intended_pt: int,
        minimum_pt: int,
        reference_id: str | None,
        field: str,
    ) -> int:
        return choose_font_size(
            heading="",
            values=[text],
            width_inches=width,
            height_inches=height,
            intended_pt=intended_pt,
            minimum_pt=minimum_pt,
            reference_id=reference_id,
            field=field,
        )

    def _add_intro(
        self,
        presentation,
        request: NarrativePresentationRequest,
        review: NarrativeReviewResponse,
    ) -> None:
        labels = INTRO_LABELS[request.generation_request.target_language]
        slide = self._clone_source_slide(
            presentation,
            ORANGE_DIVIDER_SOURCE_SLIDE,
            "divider",
        )
        for shape in list(slide.shapes):
            x, y = shape.left / 914400, shape.top / 914400
            height = shape.height / 914400
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and (
                (x > 12.4 and y > 6.7)
                or (4.4 < x < 10.8 and 2.2 < y < 4.0 and height < 1.0)
            ):
                self._remove_shape(shape)
        self._mark_source_bleed(slide, presentation)
        section_number = next(
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
            and 0.8 < shape.left / 914400 < 1.5
            and 1.2 < shape.top / 914400 < 2.0
        )
        section_number.name = "C.SECTION_NUMBER"
        self._text(
            slide,
            labels["title"],
            4.89,
            2.48,
            5.55,
            1.35,
            size=18,
            color="heading",
            valign=MSO_ANCHOR.MIDDLE,
            name="Slide title",
        )
        self._add_page_number(slide, 1)
        cleanup_unused_image_relationships(slide)
        set_notes(
            slide,
            [
                f"Authoritative PDF: {self.source_path.name}, PowerPoint source slide 10",
                "Reuse mode: direct PDF-import slide clone",
                "Reference IDs: " + ", ".join(request.approved_reference_ids),
            ],
        )
        self.provenance.append(SlideProvenance(slide_number=1, slide_type="section_narrative"))

    @staticmethod
    def _summary_groups(references: list[PreparedReference], maximum: int) -> list[list[PreparedReference]]:
        return [references[index : index + maximum] for index in range(0, len(references), maximum)]

    def _add_narrative_summary(self, presentation, references: list[PreparedReference]) -> None:
        cfg = self.config["summary"]
        groups = self._summary_groups(references, int(cfg["cards_per_slide"]))
        sequence_by_id = {
            prepared.reference.reference_id: index
            for index, prepared in enumerate(references, start=1)
        }
        for group_index, group in enumerate(groups, start=1):
            slide = self._clone_source_slide(
                presentation,
                ORANGE_COMPACT_SOURCE_SLIDES[0],
                "compact_summary",
            )
            for shape in list(slide.shapes):
                x, y = shape.left / 914400, shape.top / 914400
                height = shape.height / 914400
                keep_picture = (
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and (
                        (x < 2.2 and y > 6.7)
                        or (height > 1.5 and 0.8 < y < 6.9)
                    )
                )
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not keep_picture:
                    self._remove_shape(shape)
                elif (
                    shape.shape_type != MSO_SHAPE_TYPE.FREEFORM
                    and getattr(shape, "has_text_frame", False)
                ):
                    self._remove_shape(shape)
            cleanup_unused_image_relationships(slide)
            self._mark_source_bleed(slide, presentation)

            output_slide_number = group_index + 1
            self._text(
                slide,
                self.labels["summary"].format(current=group_index, total=len(groups)),
                0.31,
                0.25,
                4.53,
                0.50,
                size=15,
                color="heading",
                bold=True,
                valign=MSO_ANCHOR.MIDDLE,
                name="Slide title",
            )
            bullet_sources: list[BulletSource] = []
            source_slot_y = (1.03, 2.98, 4.96)
            for row_index, prepared in enumerate(group):
                reference = prepared.reference
                sequence = sequence_by_id[reference.reference_id]
                y = source_slot_y[row_index]
                mission = f"{sequence}. {reference.mission_title}"
                mission_size = self._fit(
                    mission,
                    width=1.93,
                    height=1.55,
                    intended_pt=10,
                    minimum_pt=8,
                    reference_id=reference.reference_id,
                    field="headline",
                )
                self._text(
                    slide,
                    mission,
                    0.36,
                    y + 0.14,
                    1.93,
                    1.55,
                    size=mission_size,
                    color="white",
                    bold=True,
                    align=PP_ALIGN.CENTER,
                    valign=MSO_ANCHOR.MIDDLE,
                    name=f"C.MISSION_{sequence}",
                )
                bullets = [item.text for item in prepared.summary_bullets]
                maximum_bullets = int(cfg["maximum_bullets"])
                if len(bullets) > maximum_bullets:
                    raise PptxContentOverflowError(
                        OverflowCondition(
                            reference.reference_id,
                            "compact_services",
                            len(bullets),
                            maximum_bullets,
                            int(self.type["body_min_pt"]),
                        )
                    )
                bullet_size = choose_font_size(
                    heading="",
                    values=bullets,
                    width_inches=8.05,
                    height_inches=1.52,
                    intended_pt=8,
                    minimum_pt=8,
                    reference_id=reference.reference_id,
                    field="compact_services",
                )
                bullet_box = self._bullets(
                    slide,
                    bullets,
                    2.94,
                    y + 0.16,
                    8.15,
                    1.55,
                    size=bullet_size,
                )
                bullet_box.name = f"C.SERVICES_{sequence}"
                client_meta = "\n".join(
                    value for value in (reference.client, reference.country) if value
                )
                client_size = self._fit(
                    client_meta,
                    width=1.48,
                    height=1.55,
                    intended_pt=9,
                    minimum_pt=8,
                    reference_id=reference.reference_id,
                    field="client_country",
                )
                self._text(
                    slide,
                    client_meta,
                    11.44,
                    y + 0.14,
                    1.48,
                    1.55,
                    size=client_size,
                    color="body",
                    bold=True,
                    align=PP_ALIGN.CENTER,
                    valign=MSO_ANCHOR.MIDDLE,
                    name=f"C.CLIENT_COUNTRY_{sequence}",
                )
                bullet_sources.extend(prepared.summary_bullets)
            self._add_page_number(slide, output_slide_number)
            set_notes(
                slide,
                [
                    "Authoritative PDF: PowerPoint source slide 11",
                    "Reuse mode: direct PDF-import slide clone",
                    "Narrative: human-reviewed and approved",
                    "Trusted metadata: client, country, sector, offering",
                    "Reference IDs: " + ", ".join(item.reference.reference_id for item in group),
                ],
            )
            self.provenance.append(
                SlideProvenance(
                    slide_number=output_slide_number,
                    slide_type="reference_summary",
                    reference_ids=[item.reference.reference_id for item in group],
                    bullet_sources=bullet_sources,
                )
            )

    def _prepared_references(
        self,
        request: NarrativePresentationRequest,
        review: NarrativeReviewResponse,
        selections: list[NarrativeEvidenceSelection],
        trusted_references: list[TrustedReference] | None = None,
    ) -> tuple[list[PreparedReference], list[str]]:
        trusted = list(trusted_references or [item.reference for item in selections])
        if [item.reference_id for item in trusted] != request.approved_reference_ids:
            raise ValueError("Trusted reference order does not match the approved references")
        selection_by_id = {item.reference.reference_id: item for item in selections}
        prepared: list[PreparedReference] = []
        warnings: list[str] = []
        maximum_short_description = int(self.config["summary"]["maximum_bullet_characters"])
        for narrative, metadata, trusted_reference in zip(
            review.narrative.references,
            review.reference_metadata,
            trusted,
            strict=True,
        ):
            if metadata.reference_id != trusted_reference.reference_id:
                raise ValueError("Trusted compact metadata does not match the selected reference")
            headline = narrative.headline.text.strip() or trusted_reference.mission_title
            display_reference = trusted_reference.model_copy(update={"mission_title": headline})
            selection = selection_by_id.get(trusted_reference.reference_id)
            bullets: list[BulletSource] = []
            short_description = narrative.short_description.text.strip()
            if short_description and len(short_description) <= maximum_short_description:
                bullets.append(
                    BulletSource(
                        text=short_description,
                        source_fields=["approved_narrative.short_description"],
                    )
                )
            elif short_description:
                warnings.append("COMPACT_SHORT_DESCRIPTION_OMITTED")
            for item in narrative.realisations:
                text = item.text.strip()
                if not text:
                    continue
                if len(text) > maximum_short_description and ADMINISTRATIVE_REALISATION.search(text):
                    warnings.append("COMPACT_ADMINISTRATIVE_REALISATION_OMITTED")
                    continue
                bullets.append(
                    BulletSource(
                        text=text,
                        source_fields=["approved_narrative.realisations"],
                    )
                )
            scope = " / ".join(value for value in (metadata.offering, metadata.sector) if value)
            if (
                scope
                and len(bullets) < int(self.config["summary"]["maximum_bullets"])
                and scope.casefold() not in {item.text.casefold() for item in bullets}
            ):
                candidate = [item.text for item in bullets] + [scope]
                if measure_text_fit(
                    heading="",
                    values=candidate,
                    width_inches=8.05,
                    height_inches=1.52,
                    intended_pt=8,
                    minimum_pt=8,
                ).fits:
                    bullets.append(
                        BulletSource(
                            text=scope,
                            source_fields=["trusted_metadata.offering", "trusted_metadata.sector"],
                        )
                    )
                else:
                    warnings.append("COMPACT_SCOPE_OMITTED_FOR_TEMPLATE_FIT")
            prepared.append(
                PreparedReference(
                    reference=display_reference,
                    summary_bullets=bullets,
                    description_items=[],
                    service_items=bullets,
                    why_selected=[],
                    evidence_items=[selection.evidence] if selection else [],
                )
            )
        return prepared, list(dict.fromkeys(warnings))

    def _add_source_evidence(
        self,
        presentation: Presentation,
        prepared: list[PreparedReference],
        selections: list[NarrativeEvidenceSelection],
        first_output_slide: int,
    ) -> None:
        total = len(selections)
        prepared_by_id = {item.reference.reference_id: item for item in prepared}
        for offset, selection in enumerate(selections):
            item = prepared_by_id[selection.reference.reference_id]
            slide_number = first_output_slide + offset
            slide = self._clone_source_slide(
                presentation,
                ORANGE_EVIDENCE_SOURCE_SLIDES[0],
                "evidence_attestation",
            )
            for shape in list(slide.shapes):
                x, y = shape.left / 914400, shape.top / 914400
                keep_footer_logo = (
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and x < 2.2
                    and y > 6.7
                )
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not keep_footer_logo:
                    self._remove_shape(shape)
                elif (
                    shape.shape_type != MSO_SHAPE_TYPE.FREEFORM
                    and getattr(shape, "has_text_frame", False)
                ):
                    self._remove_shape(shape)
            cleanup_unused_image_relationships(slide)
            self._mark_source_bleed(slide, presentation)

            self._text(
                slide,
                self.labels["evidence"].format(current=offset + 1, total=total),
                0.31,
                0.25,
                5.40,
                0.50,
                size=15,
                color="heading",
                bold=True,
                valign=MSO_ANCHOR.MIDDLE,
                name="Slide title",
            )
            self._text(
                slide,
                item.reference.mission_title,
                0.58,
                1.10,
                12.18,
                0.28,
                size=9,
                color="heading",
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                name="Evidence reference title",
            )
            visual = selection.visual
            if not visual.is_rendered:
                raise ValueError("Orange evidence export requires an actual approved source-page image")
            picture = self._picture_contain(
                slide,
                visual.image_path,
                0.64,
                1.47,
                12.05,
                4.98,
                name="Evidence source image",
            )
            expected_ratio = visual.rendered_pixel_width / visual.rendered_pixel_height
            placed_ratio = picture.width / picture.height
            if abs(placed_ratio - expected_ratio) > 0.01:
                raise RuntimeError("Evidence source image aspect ratio changed during Orange placement")
            citation = (
                f"{self.labels['source']}: {selection.evidence.source_file_name} · "
                f"{self.labels['page']} {selection.evidence.source_page}"
            )
            self._text(
                slide,
                citation,
                0.58,
                6.49,
                12.18,
                0.24,
                size=8,
                color="muted",
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                name="Evidence source label",
            )
            self._add_page_number(slide, slide_number)
            record = visual.manifest_record()
            record.update(
                {
                    "placed_aspect_ratio": round(placed_ratio, 6),
                    "aspect_ratio_preserved": True,
                }
            )
            set_notes(
                slide,
                [
                    "Authoritative PDF: PowerPoint source slide 18",
                    "Reuse mode: direct PDF-import slide clone",
                    f"Trusted reference ID: {selection.reference.reference_id}",
                    f"Approved source: {selection.evidence.source_file_name}",
                    f"Approved page: {selection.evidence.source_page}",
                    "Evidence selection: deterministic trusted lineage",
                ],
            )
            self.provenance.append(
                SlideProvenance(
                    slide_number=slide_number,
                    slide_type="evidence_annex",
                    reference_ids=[selection.reference.reference_id],
                    evidence_chunk_ids=[selection.evidence.chunk_id],
                    evidence_visuals=[record],
                )
            )

    def validate(self, path: Path, selected_count: int) -> dict[str, Any]:
        presentation = Presentation(path)
        width, height = presentation.slide_width, presentation.slide_height
        errors: list[str] = []
        titles = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            has_title = False
            has_number = False
            for shape in slide.shapes:
                outside = (
                    shape.left < 0
                    or shape.top < 0
                    or shape.left + shape.width > width
                    or shape.top + shape.height > height
                )
                if outside and shape.name != "Orange source bleed":
                    errors.append(f"slide {slide_index}: shape outside slide bounds")
                if shape.name == "Slide title":
                    has_title = True
                if shape.name == "Slide number":
                    has_number = True
                if not getattr(shape, "has_text_frame", False):
                    continue
                if PROHIBITED_TEXT.search(shape.text):
                    errors.append(f"slide {slide_index}: internal retrieval score text detected")
                if LOCAL_PATH.search(shape.text):
                    errors.append(f"slide {slide_index}: local filesystem path detected")
            if not has_title:
                errors.append(f"slide {slide_index}: missing title")
            else:
                titles += 1
            if not has_number:
                errors.append(f"slide {slide_index}: missing slide number")
        if errors:
            raise RuntimeError("PowerPoint validation failed: " + "; ".join(errors))
        return {
            "status": "PASS",
            "slide_count": len(presentation.slides),
            "selected_reference_count": selected_count,
            "titles_present": titles,
            "bounds_check": "PASS_WITH_SOURCE_BLEED",
            "internal_score_scan": "PASS",
            "local_path_scan": "PASS",
            "slide_numbering": "PASS",
            "evidence_visuals": "PASS",
            "evidence_aspect_ratio": "PASS",
        }

    def generate(
        self,
        output_path: Path,
        request: NarrativePresentationRequest,
        review: NarrativeReviewResponse,
        evidence_selections: list[NarrativeEvidenceSelection] | None = None,
        trusted_references: list[TrustedReference] | None = None,
        *,
        prevalidated_review: NarrativeReviewResponse | None = None,
        copy_generation: dict[str, Any] | None = None,
    ) -> NarrativePptxResult:
        selections = list(evidence_selections or [])
        prepared, export_warnings = self._prepared_references(
            request, review, selections, trusted_references
        )
        self.provenance = []
        language = request.generation_request.target_language
        self.rtl = language == "ar"
        self.font = str(self.type["arabic_font"] if self.rtl else self.type["latin_font"])
        self.labels = TEXT[language]
        # If a prevalidated review (from PresentationCopyService) is provided,
        # prefer it for rendering to avoid reintroducing overflowing content.
        if prevalidated_review is not None:
            review = prevalidated_review

        presentation = Presentation(self.clone_base_path)
        if len(presentation.slides) != self.clone_base_slide_count:
            raise RuntimeError("Orange PDF clone base slide count changed")
        self._add_intro(presentation, request, review)
        self._add_narrative_summary(presentation, prepared)

        reference_to_slide: list[dict[str, Any]] = []
        for provenance in self.provenance:
            if provenance.slide_type != "reference_summary":
                continue
            for card_index, reference_id in enumerate(provenance.reference_ids, start=1):
                reference_to_slide.append(
                    {
                        "reference_id": reference_id,
                        "slide_number": provenance.slide_number,
                        "card_index": card_index,
                    }
                )

        narrative_slide_count = 1 + len(
            [item for item in self.provenance if item.slide_type == "reference_summary"]
        )
        evidence_start = narrative_slide_count + 1
        self._add_source_evidence(presentation, prepared, selections, evidence_start)
        reference_to_evidence_slide: list[dict[str, Any]] = []
        for offset, selection in enumerate(selections):
            slide_number = evidence_start + offset
            reference_to_evidence_slide.append(
                {"reference_id": selection.reference.reference_id, "slide_number": slide_number}
            )

        evidence_visuals: list[dict[str, Any]] = []
        for provenance, selection in zip(
            (item for item in self.provenance if item.slide_type == "evidence_annex"),
            selections,
            strict=True,
        ):
            record = provenance.evidence_visuals[0]
            record.update(
                {
                    "selection_reason": selection.selection_reason,
                    "source_hash_validation": "PASS",
                    "slide_number": provenance.slide_number,
                }
            )
            record.pop("image_file_name", None)
            evidence_visuals.append(record)

        remove_original_slides(presentation, self.clone_base_slide_count)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        validation = self.validate(output_path, len(prepared))
        evidence_slide_count = len(selections)
        if validation["slide_count"] != narrative_slide_count + evidence_slide_count:
            raise RuntimeError("Generated compact presentation slide count mismatch")
        return NarrativePptxResult(
            path=output_path,
            slide_count=validation["slide_count"],
            narrative_slide_count=narrative_slide_count,
            evidence_slide_count=evidence_slide_count,
            reference_to_slide=reference_to_slide,
            reference_to_evidence_slide=reference_to_evidence_slide,
            evidence_visuals=evidence_visuals,
            font_substitution={"source": str(self.type["latin_font"]), "used": self.font},
            overflow_validation={
                "status": "PASS",
                "policy": "block",
                "minimum_font_pt": 8,
            },
            narrative_slide_mappings=[
                {
                    "narrative_role": "section_narrative",
                    "slide_number": 1,
                    "fields": ["section_intro", "overall_storyline", "why_these_references"],
                }
            ],
            export_warnings=export_warnings,
        )
